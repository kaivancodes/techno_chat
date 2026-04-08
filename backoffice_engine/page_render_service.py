"""
page_render_service.py
───────────────────────
Source viewer service. Two public functions:
  get_page_render()   — PDF page → PNG image with optional highlight (cached)
  get_source_content() — Non-PDF source → formatted text excerpt
"""

import csv as csv_module
import fitz
import hashlib
import os
import re
import textwrap
from pathlib import Path
import openpyxl
from docx import Document as DocxDoc
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageDraw, ImageFont
from django.conf import settings
from .clients import VLMClient
from .document_reader import docx_file_text
from .helpers import normalize_file_type
from .models import File
from .exceptions import PageRenderError
from .constants import PAGE_RENDER_HIGHLIGHT_MAX_LEN
from techno_chat.settings import PAGE_RENDER_DPI, PAGE_RENDER_FORMAT, logger


def _file_exists(file_obj: File) -> bool:
    try:
        return bool(file_obj.file and file_obj.file.name and os.path.exists(file_obj.file.path))
    except Exception:
        return False


def _preview_unavailable_message(file_obj: File, highlight_text: str = "", location_label: str = "") -> str:
    parts = [
        f"Preview is unavailable for {file_obj.original_filename or file_obj.file.name}.",
    ]
    if location_label:
        parts.append(f"Requested source: {location_label}.")
    parts.append("The original uploaded file could not be opened, so a page preview could not be generated.")
    return "\n".join(parts)


def get_source_fallback_content(
    file_id: int,
    file_type: str,
    page_index: int = None,
    slide_index: int = None,
    sheet_name: str = None,
    row_start: int = None,
    line_start: int = None,
    line_end: int = None,
    section_name: str = None,
    highlight_text: str = "",
) -> str:
    file_obj = File.objects.get(id=file_id)
    ft = normalize_file_type(file_type or file_obj.file_type or "", file_obj.original_filename or file_obj.file.name)

    location_bits = []
    if page_index is not None:
        location_bits.append(f"page {page_index}")
    if slide_index is not None:
        location_bits.append(f"slide {slide_index}")
    if sheet_name:
        location_bits.append(f"sheet {sheet_name}")
    if row_start is not None:
        location_bits.append(f"row {row_start}")
    if line_start is not None:
        location_bits.append(f"lines {line_start}-{line_end or line_start}")
    if section_name:
        location_bits.append(f"section {section_name}")

    location_label = ", ".join(location_bits)

    if ft in ("image", "png", "jpg", "jpeg", "webp", "svg") and _file_exists(file_obj):
        return f"[Image file]\n{_preview_unavailable_message(file_obj, highlight_text=highlight_text, location_label=location_label)}"

    return _preview_unavailable_message(file_obj, highlight_text=highlight_text, location_label=location_label)


def _load_font(size: int, bold: bool = False):
    candidates = ["DejaVuSans-Bold.ttf", "Arial Bold.ttf"] if bold else ["DejaVuSans.ttf", "Arial.ttf"]
    for candidate in candidates:
        try:
            return ImageFont.truetype(candidate, size)
        except Exception:
            continue
    return ImageFont.load_default()


def _wrap_for_width(draw: ImageDraw.ImageDraw, text: str, font, max_width: int) -> list[str]:
    wrapped = []
    for paragraph in (text or "").splitlines() or [""]:
        cleaned = paragraph.strip()
        if not cleaned:
            wrapped.append("")
            continue
        current = ""
        for word in cleaned.split():
            tentative = f"{current} {word}".strip()
            if draw.textlength(tentative, font=font) <= max_width:
                current = tentative
            else:
                if current:
                    wrapped.append(current)
                current = word
        if current:
            wrapped.append(current)
    return wrapped


def _render_text_canvas(title: str, subtitle: str, body: str, output_path: Path, width: int = 1240, height: int = 1754) -> str:
    image = Image.new("RGB", (width, height), "#faf8f2")
    draw = ImageDraw.Draw(image)
    title_font = _load_font(34, bold=True)
    subtitle_font = _load_font(20)
    body_font = _load_font(22)
    margin_x = 70
    y = 60

    draw.text((margin_x, y), title, fill="#111827", font=title_font)
    y += 60
    if subtitle:
        draw.text((margin_x, y), subtitle, fill="#475569", font=subtitle_font)
        y += 50

    draw.line((margin_x, y, width - margin_x, y), fill="#cbd5e1", width=2)
    y += 35

    lines = _wrap_for_width(draw, body, body_font, width - (margin_x * 2))
    line_height = 32
    max_lines = max(1, (height - y - 50) // line_height)
    for line in lines[:max_lines]:
        draw.text((margin_x, y), line, fill="#0f172a", font=body_font)
        y += line_height

    output_path.parent.mkdir(parents=True, exist_ok=True)
    image.save(output_path)
    return settings.MEDIA_URL + f"page_renders/{output_path.name}"


def _visual_cache_path(file_id: int, key: str) -> Path:
    cache_dir = Path(settings.MEDIA_ROOT) / "page_renders"
    cache_dir.mkdir(parents=True, exist_ok=True)
    safe_key = re.sub(r"[^a-zA-Z0-9_.-]", "_", key)
    return cache_dir / f"{file_id}_{safe_key}.{PAGE_RENDER_FORMAT}"


def _render_docx_page(file_obj: File, page_index: int) -> str:
    cache_path = _visual_cache_path(file_obj.id, f"docx_p{page_index}")
    if cache_path.exists():
        return settings.MEDIA_URL + f"page_renders/{cache_path.name}"

    segments = docx_file_text(file_obj.file.path)
    if not segments:
        raise PageRenderError(internal=f"DOCX render failed. file_id={file_obj.id} page={page_index} raw=no segments")

    segment = next((item for item in segments if item.get("page_index") == page_index), None)
    if segment is None:
        segment = min(
            segments,
            key=lambda item: abs((item.get("page_index") or 1) - page_index),
        )
    body = segment.get("text", "").strip() or "No readable content was found on this page."
    return _render_text_canvas(
        title=file_obj.original_filename or file_obj.file.name,
        subtitle=f"Document page {segment.get('page_index') or page_index}",
        body=body,
        output_path=cache_path,
    )


def _draw_boxed_text(draw, box, text, title_font, body_font, fill="#111827", outline="#cbd5e1"):
    left, top, right, bottom = box
    draw.rounded_rectangle((left, top, right, bottom), radius=12, outline=outline, width=2, fill="#ffffff")
    inner_x = left + 16
    inner_y = top + 12
    max_width = max(40, right - left - 32)
    lines = _wrap_for_width(draw, text, body_font, max_width)
    line_height = 26
    max_lines = max(1, (bottom - inner_y - 12) // line_height)
    for line in lines[:max_lines]:
        draw.text((inner_x, inner_y), line, fill=fill, font=body_font)
        inner_y += line_height


def _render_pptx_slide(file_obj: File, slide_index: int) -> str:
    cache_path = _visual_cache_path(file_obj.id, f"pptx_s{slide_index}")
    if cache_path.exists():
        return settings.MEDIA_URL + f"page_renders/{cache_path.name}"

    presentation = Presentation(file_obj.file.path)
    slide_pos = max(0, slide_index - 1)
    if slide_pos >= len(presentation.slides):
        raise PageRenderError(internal=f"PPTX render failed. file_id={file_obj.id} slide={slide_index} raw=slide out of range")

    slide = presentation.slides[slide_pos]
    width = 1280
    height = max(720, int(width * (presentation.slide_height / presentation.slide_width)))
    image = Image.new("RGB", (width, height), "#f8fafc")
    draw = ImageDraw.Draw(image)
    title_font = _load_font(28, bold=True)
    body_font = _load_font(20)
    draw.text((36, 22), f"{file_obj.original_filename or file_obj.file.name} · Slide {slide_index}", fill="#0f172a", font=title_font)

    scale_x = width / float(presentation.slide_width or 1)
    scale_y = height / float(presentation.slide_height or 1)

    for shape in slide.shapes:
        left = int(shape.left * scale_x)
        top = int(shape.top * scale_y)
        box_width = max(80, int(shape.width * scale_x))
        box_height = max(40, int(shape.height * scale_y))
        right = min(width - 20, left + box_width)
        bottom = min(height - 20, top + box_height)
        if right <= left or bottom <= top:
            continue

        if getattr(shape, "has_text_frame", False):
            text = shape.text_frame.text.strip()
            if text:
                _draw_boxed_text(draw, (left, top, right, bottom), text, title_font, body_font)
        elif getattr(shape, "has_table", False):
            cell_lines = []
            for row in shape.table.rows:
                cell_lines.append(" | ".join(cell.text.strip() for cell in row.cells if cell.text.strip()))
            if cell_lines:
                _draw_boxed_text(draw, (left, top, right, bottom), "\n".join(cell_lines), title_font, body_font)
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            draw.rounded_rectangle((left, top, right, bottom), radius=12, outline="#94a3b8", width=2, fill="#e2e8f0")
            draw.text((left + 18, top + 18), "Image", fill="#334155", font=body_font)

    cache_path.parent.mkdir(parents=True, exist_ok=True)
    image.save(cache_path)
    return settings.MEDIA_URL + f"page_renders/{cache_path.name}"


def get_visual_render(file_id: int, file_type: str, page_index: int | None = None, slide_index: int | None = None) -> str:
    file_obj = File.objects.get(id=file_id)
    if not _file_exists(file_obj):
        raise PageRenderError(internal=f"Visual render failed. file_id={file_id} raw=missing file")

    normalized = normalize_file_type(file_type or file_obj.file_type or "", file_obj.original_filename or file_obj.file.name)
    if normalized == "pdf":
        return get_page_render(file_id=file_id, page_index=page_index or 1, highlight_text="")
    if normalized in {"doc", "docx"}:
        return _render_docx_page(file_obj, page_index or 1)
    if normalized in {"ppt", "pptx"}:
        return _render_pptx_slide(file_obj, slide_index or 1)
    if normalized in {"image", "png", "jpg", "jpeg", "webp", "svg"}:
        return file_obj.file.url
    raise PageRenderError(internal=f"Visual render unsupported. file_id={file_id} type={normalized}")

def get_page_render(file_id: int, page_index: int, highlight_text: str = '') -> str:
    """
    Render one PDF page to a PNG image, draw yellow highlight if text found.

    Args:
        file_id:        File.id primary key.
        page_index:     1-based page number.
        highlight_text: Text to highlight. Empty = no highlight.

    Returns:
        URL string, e.g. /media/page_renders/42_p2_abc123def456.png

    Cache key: sha256(str(file_id) + str(page_index) + highlight_text[:50])
    Cache dir:  MEDIA_ROOT/page_renders/
    Raises:     PageRenderError
    """
    cache_key  = hashlib.sha256(
        (str(file_id) + str(page_index) + highlight_text[:50]).encode()
    ).hexdigest()[:12]
    filename   = f'{file_id}_p{page_index}_{cache_key}.{PAGE_RENDER_FORMAT}'
    cache_dir  = Path(settings.MEDIA_ROOT) / 'page_renders'
    cache_dir.mkdir(parents=True, exist_ok=True)
    cache_path = cache_dir / filename

    if cache_path.exists():
        logger.info('get_page_render | cache hit | %s', filename)
        return settings.MEDIA_URL + 'page_renders/' + filename

    try:
        file_obj  = File.objects.get(id=file_id)
        if not _file_exists(file_obj):
            raise FileNotFoundError(file_obj.file.path)
        file_path = file_obj.file.path
        doc       = fitz.open(file_path)
        page      = doc[page_index - 1]

        # Draw highlight if text provided
        if highlight_text:
            search_str = highlight_text[:PAGE_RENDER_HIGHLIGHT_MAX_LEN]
            rects = page.search_for(search_str)
            for rect in rects:
                page.draw_rect(
                    rect,
                    color=(1, 1, 0),
                    fill=(1, 1, 0.5),
                    fill_opacity=0.35,
                    overlay=True,
                )

        pix = page.get_pixmap(dpi=PAGE_RENDER_DPI)
        pix.save(str(cache_path))
        doc.close()
        logger.info('get_page_render | rendered | %s', filename)
        return settings.MEDIA_URL + 'page_renders/' + filename

    except Exception as exc:
        raise PageRenderError(internal=f'PDF render failed. file_id={file_id} page={page_index} raw={exc}')

def get_source_content(
    file_id: int,
    file_type: str,
    page_index: int   = None,
    slide_index: int  = None,
    sheet_name: str   = None,
    row_start: int    = None,
    line_start: int   = None,
    line_end: int     = None,
    section_name: str = None,
    highlight_text: str = '',
) -> str:
    """
    For non-PDF files: extract and return a readable text excerpt from the
    location that produced the answer. Used by the source viewer modal for
    DOCX, PPTX, XLSX, CSV, TXT, MD, and image files.

    Returns a formatted text string (max ~1500 chars) describing the source.
    Raises PageRenderError on failure.
    """
    try:
        file_obj  = File.objects.get(id=file_id)
        if not _file_exists(file_obj):
            return get_source_fallback_content(
                file_id=file_id,
                file_type=file_type,
                page_index=page_index,
                slide_index=slide_index,
                sheet_name=sheet_name,
                row_start=row_start,
                line_start=line_start,
                line_end=line_end,
                section_name=section_name,
                highlight_text=highlight_text,
            )
        file_path = file_obj.file.path
        ft        = normalize_file_type(file_type or file_obj.file_type or "", file_obj.original_filename or file_obj.file.name)

        # Image files: return VLM description (re-run if not cached)
        if ft in ('image', 'png', 'jpg', 'jpeg', 'webp', 'svg'):
            try:
                text = VLMClient().describe_svg_file(file_path) if ft == "svg" else VLMClient().describe_image_file(file_path)
                return f'[Image file description]\n{text[:1500]}'
            except Exception:
                return get_source_fallback_content(
                    file_id=file_id,
                    file_type=file_type,
                    page_index=page_index,
                    slide_index=slide_index,
                    sheet_name=sheet_name,
                    row_start=row_start,
                    line_start=line_start,
                    line_end=line_end,
                    section_name=section_name,
                    highlight_text=highlight_text,
                )

        # PPTX: extract text from the relevant slide
        if ft in ('ppt', 'pptx'):
            prs  = Presentation(file_path)
            target_slide = slide_index or 1
            idx  = target_slide - 1
            if 0 <= idx < len(prs.slides):
                parts = [f'[Slide {target_slide}]']
                for shape in prs.slides[idx].shapes:
                    if shape.has_text_frame:
                        parts.append(shape.text_frame.text.strip())
                    elif getattr(shape, "has_table", False):
                        for row in shape.table.rows:
                            parts.append(" | ".join(cell.text.strip() for cell in row.cells))
                return '\n'.join(p for p in parts if p)[:1500]
            return f'Slide {slide_index} not found in presentation.'

        # XLSX: extract relevant sheet and row range
        if ft in ('excel', 'xlsx', 'xls'):
            wb = openpyxl.load_workbook(file_path, data_only=True)
            target_sheet = sheet_name or (wb.sheetnames[0] if wb.sheetnames else None)
            if not target_sheet:
                return get_source_fallback_content(
                    file_id=file_id,
                    file_type=file_type,
                    page_index=page_index,
                    slide_index=slide_index,
                    sheet_name=sheet_name,
                    row_start=row_start,
                    line_start=line_start,
                    line_end=line_end,
                    section_name=section_name,
                    highlight_text=highlight_text,
                )
            if target_sheet in wb.sheetnames:
                ws   = wb[target_sheet]
                rows = []
                start = (row_start or 1) - 1
                for i, row in enumerate(ws.iter_rows(values_only=True)):
                    if i < start:
                        continue
                    rows.append(' | '.join('' if c is None else str(c) for c in row))
                    if len(rows) >= 60:
                        break
                return f'[Sheet: {target_sheet}]\n' + '\n'.join(rows)[:1500]
            return f'Sheet "{target_sheet}" not found.'

        # CSV: extract relevant row range
        if ft == 'csv':
            lines = []
            with open(file_path, 'r', encoding='utf-8', errors='replace') as fh:
                for i, row in enumerate(csv_module.reader(fh), 1):
                    if i >= (row_start or 1):
                        lines.append(' | '.join(str(c) for c in row))
                    if len(lines) >= 60:
                        break
            return f'[CSV rows from {row_start or 1}]\n' + '\n'.join(lines)[:1500]

        # TXT / MD / DOCX: read and return a window of text
        if ft == 'txt':
            with open(file_path, 'r', encoding='utf-8', errors='replace') as fh:
                lines = fh.read().splitlines()
            start = max(1, line_start or 1)
            end = max(start, line_end or min(start + 40, len(lines)))
            excerpt = []
            for index in range(start, min(end, len(lines)) + 1):
                excerpt.append(lines[index - 1])
            return '\n'.join(excerpt)[:1500]

        if ft == 'md':
            with open(file_path, 'r', encoding='utf-8', errors='replace') as fh:
                lines = fh.read().splitlines()
            if section_name:
                return _extract_md_named_section(lines, section_name)[:1500]
            start = max(1, line_start or 1)
            end = max(start, line_end or min(start + 40, len(lines)))
            excerpt = []
            for index in range(start, min(end, len(lines)) + 1):
                excerpt.append(lines[index - 1])
            return '\n'.join(excerpt)[:1500]

        if ft in ('doc', 'docx'):
            doc   = DocxDoc(file_path)
            paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
            start = (page_index - 1) * 5 if page_index else 0
            return '\n'.join(paras[start:start + 30])[:1500]

        return get_source_fallback_content(
            file_id=file_id,
            file_type=file_type,
            page_index=page_index,
            slide_index=slide_index,
            sheet_name=sheet_name,
            row_start=row_start,
            line_start=line_start,
            line_end=line_end,
            section_name=section_name,
            highlight_text=highlight_text,
        )

    except Exception as exc:
        raise PageRenderError(
            internal=f'get_source_content failed. file_id={file_id} ft={file_type} raw={exc}'
        )


def _extract_md_named_section(lines: list[str], section_name: str) -> str:
    heading_pattern = re.compile(r"^\s{0,3}#{1,6}\s+(.+?)\s*$")
    collecting = False
    extracted_lines = []

    for line in lines:
        match = heading_pattern.match(line)
        if match:
            heading = match.group(1).strip()
            if collecting:
                break
            if heading == section_name:
                collecting = True
        if collecting:
            extracted_lines.append(line)

    if extracted_lines:
        return '\n'.join(extracted_lines)
    return '\n'.join(lines)
