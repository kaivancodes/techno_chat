"""
structured_file_service.py
──────────────────────────
Deterministic answering for structured documents.

Supports:
  - TXT / MD:  cricket-stats player data parser (existing)
  - CSV / XLSX: generic tabular statistics engine (pandas)
  - PDF:        page-count and outline structure queries
  - DOCX:       heading/section structure queries
  - Images:     description-based queries
"""

from __future__ import annotations

import os
import re

import fitz                          # PyMuPDF  (PDF stats)
import pandas as pd                  # CSV / Excel stats
from docx import Document as DocxDoc # DOCX stats
from pptx import Presentation        # PPTX stats

from .document_reader import extract_file_text
from .models import File
from .query_service import is_list_request, is_summary_request


_COUNTRY_NAMES = (
    "India",
    "South Africa",
    "Australia",
    "West Indies",
    "Sri Lanka",
    "New Zealand",
    "England",
    "Pakistan",
    "Bangladesh",
    "Afghanistan",
)

_COUNTRY_ALIASES = {
    "indian": "India",
    "indians": "India",
    "south african": "South Africa",
    "south africans": "South Africa",
    "protea": "South Africa",
    "proteas": "South Africa",
    "australian": "Australia",
    "australians": "Australia",
    "west indian": "West Indies",
    "west indies": "West Indies",
    "west indians": "West Indies",
    "sri lankan": "Sri Lanka",
    "sri lankans": "Sri Lanka",
    "kiwi": "New Zealand",
    "kiwis": "New Zealand",
    "new zealand": "New Zealand",
    "english": "England",
    "england": "England",
    "pakistani": "Pakistan",
    "bangladeshi": "Bangladesh",
    "afghan": "Afghanistan",
}

_SECTION_HEADING_PATTERN = re.compile(r"^\s{0,3}#{1,6}\s+(.+?)\s*$")
_RUNS_PATTERN = re.compile(r"(\d[\d,]*)\s+runs\b", re.IGNORECASE)
_PLAYER_QUERY_PATTERN = re.compile(
    r"\b(?:who is|who's|tell me about|about|describe|explain)\s+([A-Z][A-Za-z'.-]*(?:\s+[A-Z][A-Za-z'.-]*){0,4})\b"
)
_NUMBER_TOKEN_PATTERN = re.compile(r"(\d[\d,]*)(\+)?")

_WORD_NUMBERS = {
    "zero": 0,
    "one": 1,
    "two": 2,
    "three": 3,
    "four": 4,
    "five": 5,
    "six": 6,
    "seven": 7,
    "eight": 8,
    "nine": 9,
    "ten": 10,
}

_STAT_LABELS = {
    "runs": "Runs",
    "batting_average": "Batting Average",
    "strike_rate": "Strike Rate",
    "centuries": "Centuries",
    "fifties": "Fifties",
    "highest_score": "Highest Score",
    "fours": "Fours",
    "sixes": "Sixes",
    "wickets": "Wickets",
    "bowling_average": "Bowling Average",
    "economy_rate": "Economy Rate",
    "bowling_strike_rate": "Bowling Strike Rate",
    "four_wicket_hauls": "4-Wicket Hauls",
    "five_wicket_hauls": "5-Wicket Hauls",
    "catches": "Catches",
    "stumpings": "Stumpings",
    "runouts": "Runouts",
}


def try_build_structured_answer(query: str, file_ids: list[int]) -> dict | None:
    normalized_query = re.sub(r"\s+", " ", (query or "").strip())
    if not normalized_query or not file_ids:
        return None

    files = list(File.objects.filter(id__in=file_ids).order_by("id"))
    for file_obj in files:
        file_type = (file_obj.file_type or "").lower()

        # ── 1. Existing cricket-stats parser (txt / md only) ──────────────────
        if file_type in {"txt", "md"}:
            structured_doc = _parse_structured_file(file_obj)
            if structured_doc:
                answer = _answer_from_structured_doc(normalized_query, structured_doc)
                if answer:
                    return answer
        elif file_type in {"power", "ppt", "pptx"}:
            structured_doc = _parse_segmented_structured_file(file_obj)
            if structured_doc:
                answer = _answer_from_structured_doc(normalized_query, structured_doc)
                if answer:
                    return answer

        # ── 2. Generic CSV / Excel tabular analysis ───────────────────────────
        elif file_type in {"csv"}:
            answer = _answer_csv_query(normalized_query, file_obj)
            if answer:
                return answer

        elif file_type in {"excel", "xlsx", "xls"}:
            answer = _answer_excel_query(normalized_query, file_obj)
            if answer:
                return answer

        elif file_type in {"power", "pptx", "ppt"}:
            answer = _answer_pptx_query(normalized_query, file_obj)
            if answer:
                return answer

        # ── 3. PDF structure queries ──────────────────────────────────────────
        elif file_type == "pdf":
            answer = _answer_pdf_query(normalized_query, file_obj)
            if answer:
                return answer

        # ── 4. DOCX structure queries ─────────────────────────────────────────
        elif file_type in {"doc", "docx"}:
            answer = _answer_docx_query(normalized_query, file_obj)
            if answer:
                return answer

        # ── 5. Image description queries ──────────────────────────────────────
        elif file_type in {"image", "png", "jpg", "jpeg", "webp", "svg"}:
            answer = _answer_image_query(normalized_query, file_obj)
            if answer:
                return answer

    return None


def _parse_structured_file(file_obj: File) -> dict | None:
    file_type = (file_obj.file_type or "").lower()
    if file_type not in {"txt", "md"}:
        return None

    try:
        with open(file_obj.file.path, "r", encoding="utf-8", errors="replace") as file_handle:
            raw_lines = file_handle.read().splitlines()
    except OSError:
        return None

    cleaned_lines = []
    for line_number, raw_line in enumerate(raw_lines, start=1):
        line = _clean_line(raw_line)
        cleaned_lines.append({
            "line_number": line_number,
            "text": line,
            "raw_text": raw_line.rstrip("\n"),
        })

    players = []
    glossary = {}
    current_country = ""
    current_section = "Section 1"
    current_player_name = ""
    current_player_lines = []
    current_player_start = None
    current_player_section = current_section
    current_player_country = ""

    index = 0
    while index < len(cleaned_lines):
        item = cleaned_lines[index]
        line_text = item["text"]

        heading_match = _SECTION_HEADING_PATTERN.match(item["raw_text"])
        if heading_match:
            current_section = heading_match.group(1).strip() or current_section
            country_name = _extract_country_heading(current_section)
            if country_name:
                current_country = country_name
            index += 1
            continue

        country_name = _extract_country_heading(line_text)
        if country_name:
            _flush_player(
                players,
                file_obj,
                current_player_name,
                current_player_lines,
                current_player_start,
                item["line_number"] - 1,
                current_player_section,
                current_player_country,
            )
            current_player_name = ""
            current_player_lines = []
            current_player_start = None
            current_country = country_name
            current_section = country_name
            index += 1
            continue

        if not current_country and not current_player_name:
            glossary_entry = _extract_glossary_entry(cleaned_lines, index)
            if glossary_entry:
                glossary[glossary_entry["term"]] = glossary_entry
                index = glossary_entry["next_index"]
                continue

        if _is_player_heading(line_text):
            _flush_player(
                players,
                file_obj,
                current_player_name,
                current_player_lines,
                current_player_start,
                item["line_number"] - 1,
                current_player_section,
                current_player_country,
            )
            current_player_name = line_text
            current_player_lines = []
            current_player_start = item["line_number"]
            current_player_section = current_section
            current_player_country = current_country
            index += 1
            continue

        if current_player_name:
            if line_text:
                current_player_lines.append((item["line_number"], line_text))
            index += 1
            continue

        index += 1

    _flush_player(
        players,
        file_obj,
        current_player_name,
        current_player_lines,
        current_player_start,
        cleaned_lines[-1]["line_number"] if cleaned_lines else 1,
        current_player_section,
        current_player_country,
    )

    if not players and not glossary:
        return None

    return {
        "file_id": file_obj.id,
        "file_name": file_obj.original_filename or file_obj.file.name,
        "file_type": file_type,
        "players": players,
        "glossary": glossary,
    }


def _parse_segmented_structured_file(file_obj: File) -> dict | None:
    try:
        segments = extract_file_text(file_obj.file.path)
    except Exception:
        return None

    if not segments:
        return None

    players = []
    glossary = {}
    current_country = ""
    current_section = "Section 1"
    current_player_name = ""
    current_player_lines = []
    current_player_start = None
    current_player_section = current_section
    current_player_country = ""
    current_player_location = {}

    for segment in segments:
        location = {
            "page_index": segment.get("page_index"),
            "slide_index": segment.get("slide_index"),
            "sheet_name": segment.get("sheet_name"),
        }
        segment_text = segment.get("text", "")
        if not segment_text:
            continue

        cleaned_lines = []
        for line_number, raw_line in enumerate(segment_text.splitlines(), start=1):
            line = _clean_line(raw_line)
            if not line or line.startswith("[Slide ") or line.startswith("PowerPoint Presentation Statistics"):
                continue
            cleaned_lines.append({
                "line_number": line_number,
                "text": line,
                "raw_text": raw_line.rstrip("\n"),
            })

        index = 0
        while index < len(cleaned_lines):
            item = cleaned_lines[index]
            line_text = item["text"]

            heading_match = _SECTION_HEADING_PATTERN.match(item["raw_text"])
            if heading_match:
                current_section = heading_match.group(1).strip() or current_section
                country_name = _extract_country_heading(current_section)
                if country_name:
                    current_country = country_name
                index += 1
                continue

            country_name = _extract_country_heading(line_text)
            if country_name:
                _flush_player(
                    players,
                    file_obj,
                    current_player_name,
                    current_player_lines,
                    current_player_start,
                    item["line_number"] - 1,
                    current_player_section,
                    current_player_country,
                    source_location=current_player_location,
                )
                current_player_name = ""
                current_player_lines = []
                current_player_start = None
                current_country = country_name
                current_section = country_name
                current_player_location = {}
                index += 1
                continue

            if not current_country and not current_player_name:
                glossary_entry = _extract_glossary_entry(cleaned_lines, index)
                if glossary_entry:
                    glossary[glossary_entry["term"]] = glossary_entry
                    index = glossary_entry["next_index"]
                    continue

            if _is_player_heading(line_text):
                _flush_player(
                    players,
                    file_obj,
                    current_player_name,
                    current_player_lines,
                    current_player_start,
                    item["line_number"] - 1,
                    current_player_section,
                    current_player_country,
                    source_location=current_player_location,
                )
                current_player_name = line_text
                current_player_lines = []
                current_player_start = item["line_number"]
                current_player_section = current_section
                current_player_country = current_country
                current_player_location = dict(location)
                index += 1
                continue

            if current_player_name:
                current_player_lines.append((item["line_number"], line_text))
            index += 1

    _flush_player(
        players,
        file_obj,
        current_player_name,
        current_player_lines,
        current_player_start,
        current_player_lines[-1][0] if current_player_lines else current_player_start,
        current_player_section,
        current_player_country,
        source_location=current_player_location,
    )

    if not players and not glossary:
        return None

    return {
        "file_id": file_obj.id,
        "file_name": file_obj.original_filename or file_obj.file.name,
        "file_type": (file_obj.file_type or "").lower(),
        "players": players,
        "glossary": glossary,
    }


def _answer_from_structured_doc(query: str, structured_doc: dict) -> dict | None:
    lowered_query = query.lower()
    player_match = _match_player_query(query, structured_doc["players"])
    matched_players = _extract_query_players(query, structured_doc["players"])
    country_name = _match_country_in_query(query)

    if player_match:
        if _is_stats_only_query(lowered_query):
            return _build_player_stats_only_answer(structured_doc, player_match)

        requested_fields = _detect_requested_fields(query)
        if len(matched_players) > 1 and requested_fields:
            return _build_multi_player_field_answer(structured_doc, matched_players, requested_fields)
        if requested_fields:
            return _build_single_player_field_answer(structured_doc, player_match, requested_fields)

    glossary_match = _match_glossary_query(lowered_query, structured_doc["glossary"]) if _is_definition_query(lowered_query) else None
    if glossary_match:
        return {
            "answer": glossary_match["definition"],
            "sources": [_build_glossary_source(structured_doc, glossary_match)],
        }

    if _is_country_count_query(lowered_query):
        return _build_country_count_answer(structured_doc)

    if _is_country_list_query(lowered_query):
        return _build_country_list_answer(structured_doc)

    if "how many" in lowered_query and "player" in lowered_query:
        count = len(structured_doc["players"])
        if count:
            player_word = "player" if count == 1 else "players"
            return {
                "answer": f"{count} cricket {player_word} are mentioned in the file.",
                "sources": [_build_range_source(structured_doc, structured_doc["players"])],
            }

    if is_list_request(query) and country_name:
        return _build_country_players_answer(structured_doc, country_name)

    if is_list_request(query) and "player" in lowered_query:
        return _build_all_players_list_answer(structured_doc)

    if "captain" in lowered_query:
        captain_players = []
        for player in structured_doc["players"]:
            if "captain" in player["description"].lower():
                captain_players.append(player)
        if captain_players:
            answer_lines = []
            for player in captain_players:
                answer_lines.append(player["name"])
            return {
                "answer": "\n".join(answer_lines),
                "sources": [_build_range_source(structured_doc, captain_players)],
            }

    extreme_request = _detect_extreme_stat_request(lowered_query)
    if extreme_request:
        stat_key, direction = extreme_request
        return _build_extreme_stat_answer(structured_doc, stat_key, direction)

    requested_fields = _detect_requested_fields(query)
    if matched_players and requested_fields:
        return _build_multi_player_field_answer(structured_doc, matched_players, requested_fields)

    if player_match:
        return _build_player_description_answer(structured_doc, player_match, lowered_query)

    if country_name and lowered_query.startswith("what about"):
        return _build_country_context_answer(structured_doc, country_name)

    if is_summary_request(query):
        return _build_summary_answer(structured_doc)

    return None


def _build_summary_answer(structured_doc: dict) -> dict | None:
    players = structured_doc["players"]
    if not players:
        return None

    countries = []
    players_by_country = {}
    for player in players:
        country = player["country"] or "Unknown"
        if country not in countries:
            countries.append(country)
        players_by_country.setdefault(country, []).append(player["name"])

    lines = []
    lines.append(
        f"The file mentions {len(players)} cricketers from {len(countries)} countries: {', '.join(countries)}."
    )
    for country in countries:
        lines.append(f"{country}: {', '.join(players_by_country[country])}")

    return {
        "answer": "\n".join(lines),
        "sources": [_build_range_source(structured_doc, players)],
    }


def _build_country_count_answer(structured_doc: dict) -> dict | None:
    countries = _ordered_countries(structured_doc["players"])
    if not countries:
        return None

    return {
        "answer": str(len(countries)),
        "sources": [_build_range_source(structured_doc, structured_doc["players"])],
    }


def _build_country_list_answer(structured_doc: dict) -> dict | None:
    countries = _ordered_countries(structured_doc["players"])
    if not countries:
        return None

    return {
        "answer": "\n".join(countries),
        "sources": [_build_range_source(structured_doc, structured_doc["players"])],
    }


def _build_country_players_answer(structured_doc: dict, country_name: str) -> dict | None:
    country_players = []
    for player in structured_doc["players"]:
        if player["country"].lower() == country_name.lower():
            country_players.append(player)

    if not country_players:
        return None

    answer_lines = []
    for player in country_players:
        answer_lines.append(player["name"])

    return {
        "answer": "\n".join(answer_lines),
        "sources": [_build_range_source(structured_doc, country_players)],
    }


def _build_all_players_list_answer(structured_doc: dict) -> dict | None:
    players = structured_doc["players"]
    if not players:
        return None

    answer_lines = []
    for player in players:
        answer_lines.append(player["name"])

    return {
        "answer": "\n".join(answer_lines),
        "sources": [_build_range_source(structured_doc, players)],
    }


def _build_player_description_answer(structured_doc: dict, player: dict, lowered_query: str) -> dict:
    if _is_short_player_overview_query(lowered_query):
        description = _build_player_short_summary(player)
    else:
        description = player["description"]
        if lowered_query.startswith("who is"):
            description = _first_sentence(description)

    return {
        "answer": description,
        "sources": [_build_source(structured_doc, player)],
    }


def _is_short_player_overview_query(lowered_query: str) -> bool:
    return any(
        phrase in lowered_query
        for phrase in ("tell me about", "who is", "who's", "about ", "describe ")
    ) and not any(
        phrase in lowered_query
        for phrase in ("exact", "verbatim", "word for word", "word to word", "full stat", "full stats", "all stats")
    )


def _build_player_short_summary(player: dict) -> str:
    stats = player.get("stats", {})
    name = player["name"]
    country = player.get("country") or "their team"
    role_text = _extract_player_role(player.get("description", ""))

    opening = f"{name} is a notable ODI cricketer for {country}."
    if role_text:
        opening = f"{name} is a notable ODI {role_text} for {country}."

    details = []
    if stats.get("runs") is not None:
        batting_bits = [f"he scored {_format_stat_value('runs', stats['runs'])} runs"]
        if stats.get("batting_average") is not None:
            batting_bits.append(f"at a batting average of {_format_stat_value('batting_average', stats['batting_average'])}")
        if stats.get("strike_rate") is not None:
            batting_bits.append(f"with a batting strike rate of {_format_stat_value('strike_rate', stats['strike_rate'])}")
        if stats.get("centuries") is not None or stats.get("fifties") is not None:
            centuries = _format_stat_value("centuries", stats.get("centuries", 0))
            fifties = _format_stat_value("fifties", stats.get("fifties", 0))
            batting_bits.append(f"including {centuries} centuries and {fifties} fifties")
        details.append(", ".join(batting_bits).replace(", with a strike rate", " and a strike rate"))

    if stats.get("wickets") is not None:
        bowling_bits = [f"with the ball he took {_format_stat_value('wickets', stats['wickets'])} wickets"]
        if stats.get("bowling_average") is not None:
            bowling_bits.append(f"at a bowling average of {_format_stat_value('bowling_average', stats['bowling_average'])}")
        if stats.get("economy_rate") is not None:
            bowling_bits.append(f"and an economy of {_format_stat_value('economy_rate', stats['economy_rate'])}")
        details.append(", ".join(bowling_bits))

    fielding_parts = []
    if stats.get("catches") is not None:
        fielding_parts.append(f"{_format_stat_value('catches', stats['catches'])} catches")
    if stats.get("stumpings") is not None:
        fielding_parts.append(f"{_format_stat_value('stumpings', stats['stumpings'])} stumpings")
    if stats.get("runouts") is not None:
        fielding_parts.append(f"{_format_stat_value('runouts', stats['runouts'])} runouts")
    if fielding_parts:
        details.append("in the field he added " + ", ".join(fielding_parts))

    if not details:
        description_hint = _paraphrase_description_hint(player.get("description", ""), name)
        if description_hint:
            return f"{opening} {description_hint}"
        return opening

    answer = f"{opening} " + ". ".join(part[0].upper() + part[1:] if part else part for part in details) + "."
    return re.sub(r"\s+", " ", answer).strip()


def _extract_player_role(description: str) -> str:
    lowered = (description or "").lower()
    if "wicketkeeper" in lowered and "captain" in lowered:
        return "wicketkeeper-captain"
    if "all-rounder" in lowered:
        return "all-rounder"
    if "wicketkeeper" in lowered:
        return "wicketkeeper"
    if "captain" in lowered:
        return "captain"
    if "bowler" in lowered or "spinner" in lowered or "fast bowler" in lowered or "seam bowler" in lowered:
        return "bowler"
    if "batsman" in lowered or "batter" in lowered:
        return "batsman"
    return ""


def _paraphrase_description_hint(description: str, player_name: str) -> str:
    sentence = _first_sentence(description)
    if not sentence:
        return ""

    hint = sentence
    hint = re.sub(rf"^{re.escape(player_name)}\s+(is|was)\s+", "", hint, flags=re.IGNORECASE)
    hint = hint.rstrip(". ")
    if not hint:
        return ""
    return f"The file describes {player_name} as {hint.lower()}."


def _build_player_stats_only_answer(structured_doc: dict, player: dict) -> dict:
    answer_lines = []
    for field in _ordered_full_stat_fields():
        value = player["stats"].get(field)
        if value is None:
            continue
        answer_lines.append(f"{_STAT_LABELS[field]}: {_format_stat_value(field, value)}")

    return {
        "answer": "\n".join(answer_lines),
        "sources": [_build_source(structured_doc, player)],
    }


def _build_single_player_field_answer(structured_doc: dict, player: dict, requested_fields: list[str]) -> dict | None:
    answer_lines = []
    for field in requested_fields:
        value = player["stats"].get(field)
        if value is None:
            continue
        answer_lines.append(f"{player['name']}: {_format_stat_value(field, value)} {_field_unit(field)}".strip())

    if not answer_lines:
        return None

    return {
        "answer": "\n".join(answer_lines),
        "sources": [_build_source(structured_doc, player)],
    }


def _build_multi_player_field_answer(structured_doc: dict, players: list[dict], requested_fields: list[str]) -> dict | None:
    answer_lines = []
    for player in players:
        field_parts = []
        for field in requested_fields:
            value = player["stats"].get(field)
            if value is None:
                continue
            field_parts.append(f"{_STAT_LABELS[field]}: {_format_stat_value(field, value)}")
        if field_parts:
            answer_lines.append(f"{player['name']} — " + ", ".join(field_parts))

    if not answer_lines:
        return None

    return {
        "answer": "\n".join(answer_lines),
        "sources": [_build_range_source(structured_doc, players)],
    }


def _build_extreme_stat_answer(structured_doc: dict, stat_key: str, direction: str) -> dict | None:
    candidates = []
    for player in structured_doc["players"]:
        value = player["stats"].get(stat_key)
        if value is None:
            continue
        candidates.append(player)

    if not candidates:
        return None

    reverse = direction == "highest"
    candidates.sort(key=lambda item: item["stats"][stat_key], reverse=reverse)
    target_value = candidates[0]["stats"][stat_key]
    winners = []
    for player in candidates:
        if player["stats"][stat_key] == target_value:
            winners.append(player)

    label = _STAT_LABELS.get(stat_key, stat_key.replace("_", " ").title())
    direction_text = "highest" if reverse else "lowest"
    value_text = f"{_format_stat_value(stat_key, target_value)} {_field_unit(stat_key)}".strip()

    if len(winners) == 1:
        winner = winners[0]
        return {
            "answer": f"{winner['name']} has the {direction_text} {label.lower()} with {value_text}.",
            "sources": [_build_source(structured_doc, winner)],
        }

    names = []
    for winner in winners:
        names.append(winner["name"])
    return {
        "answer": f"{', '.join(names)} share the {direction_text} {label.lower()} with {value_text} each.",
        "sources": [_build_range_source(structured_doc, winners)],
    }


def _build_country_context_answer(structured_doc: dict, country_name: str) -> dict | None:
    country_players = []
    for player in structured_doc["players"]:
        if player["country"].lower() == country_name.lower():
            country_players.append(player)

    if not country_players:
        return None

    names = []
    for player in country_players:
        names.append(player["name"])

    return {
        "answer": f"{country_name} is mentioned in the file. Players: {', '.join(names)}.",
        "sources": [_build_range_source(structured_doc, country_players)],
    }


def _flush_player(
    players: list[dict],
    file_obj: File,
    name: str,
    description_lines: list[tuple[int, str]],
    start_line: int | None,
    end_line: int | None,
    section_name: str,
    country: str,
    source_location: dict | None = None,
) -> None:
    if not name:
        return

    description_parts = []
    line_numbers = []
    for line_number, text in description_lines:
        if text:
            description_parts.append(text)
            line_numbers.append(line_number)

    description = re.sub(r"\s+", " ", " ".join(description_parts)).strip()
    if not description:
        return

    cleaned_description = _sanitize_description(description)
    runs = _extract_runs(cleaned_description)
    source_location = source_location or {}
    players.append({
        "file_id": file_obj.id,
        "file_name": file_obj.original_filename or file_obj.file.name,
        "file_type": (file_obj.file_type or "").lower(),
        "name": name,
        "country": country or "",
        "section_name": section_name or "Section 1",
        "line_start": start_line or (line_numbers[0] if line_numbers else 1),
        "line_end": line_numbers[-1] if line_numbers else (end_line or start_line or 1),
        "page_index": source_location.get("page_index"),
        "slide_index": source_location.get("slide_index"),
        "sheet_name": source_location.get("sheet_name"),
        "description": cleaned_description,
        "runs": runs,
        "stats": _extract_player_stats(cleaned_description, name),
    })


def _extract_glossary_entry(cleaned_lines: list[dict], index: int) -> dict | None:
    current_text = cleaned_lines[index]["text"]
    if not current_text or len(current_text.split()) > 5:
        return None
    if _extract_country_heading(current_text):
        return None

    next_index = index + 1
    definition_lines = []
    while next_index < len(cleaned_lines):
        candidate = cleaned_lines[next_index]
        candidate_text = candidate["text"]
        if not candidate_text:
            if definition_lines:
                break
            next_index += 1
            continue
        if _extract_country_heading(candidate_text) or _is_player_heading(candidate_text):
            break
        definition_lines.append(candidate)
        next_index += 1

    if not definition_lines:
        return None

    definition_parts = []
    for item in definition_lines:
        definition_parts.append(item["text"])

    return {
        "term": current_text.lower(),
        "title": current_text,
        "definition": re.sub(r"\s+", " ", " ".join(definition_parts)).strip(),
        "line_start": cleaned_lines[index]["line_number"],
        "line_end": definition_lines[-1]["line_number"],
        "next_index": next_index,
    }


def _match_player_query(query: str, players: list[dict]) -> dict | None:
    explicit_match = _PLAYER_QUERY_PATTERN.search(query)
    if explicit_match:
        target_name = _normalize_name(explicit_match.group(1).strip())
        for player in players:
            if _normalize_name(player["name"]) == target_name:
                return player

    normalized_query = re.sub(r"[^a-z\s]", " ", query.lower())
    normalized_query = re.sub(r"\s+", " ", normalized_query).strip()
    normalized_compact_query = _normalize_name(query)
    if not normalized_query:
        return None

    for player in players:
        if player["name"].lower() == normalized_query:
            return player

    for player in players:
        if player["name"].lower() in normalized_query:
            return player
        if _normalize_name(player["name"]) in normalized_compact_query:
            return player
    return None


def _match_glossary_query(lowered_query: str, glossary: dict) -> dict | None:
    for term, entry in glossary.items():
        if term in lowered_query:
            return entry
        title = (entry.get("title") or "").lower()
        normalized_title = re.sub(r"[^a-z0-9\s]", " ", title)
        normalized_title = re.sub(r"\s+", " ", normalized_title).strip()
        if normalized_title and normalized_title in lowered_query:
            return entry
        base_title = re.sub(r"\([^)]*\)", "", title)
        base_title = re.sub(r"[^a-z0-9\s]", " ", base_title)
        base_title = re.sub(r"\s+", " ", base_title).strip()
        if base_title and base_title in lowered_query:
            return entry
        for token in [item for item in re.split(r"[^a-z0-9]+", title) if len(item) > 3]:
            if token in lowered_query:
                return entry
    return None


def _build_source(structured_doc: dict, player: dict) -> dict:
    source = {
        "file_name": structured_doc["file_name"],
        "file_type": structured_doc["file_type"],
        "file_id": structured_doc["file_id"],
        "highlight_text": player["name"],
    }
    if structured_doc["file_type"] == "md" and player.get("section_name"):
        source["section_name"] = player["section_name"]
    if player.get("line_start") is not None:
        source["line_start"] = player["line_start"]
    if player.get("line_end") is not None:
        source["line_end"] = player["line_end"]
    if player.get("page_index") is not None:
        source["page_index"] = player["page_index"]
    if player.get("slide_index") is not None:
        source["slide_index"] = player["slide_index"]
    if player.get("sheet_name"):
        source["sheet_name"] = player["sheet_name"]
    return source


def _build_glossary_source(structured_doc: dict, glossary_entry: dict) -> dict:
    return {
        "file_name": structured_doc["file_name"],
        "file_type": structured_doc["file_type"],
        "file_id": structured_doc["file_id"],
        "line_start": glossary_entry["line_start"],
        "line_end": glossary_entry["line_end"],
        "highlight_text": glossary_entry["title"],
    }


def _build_range_source(structured_doc: dict, items: list[dict]) -> dict:
    line_start = None
    line_end = None
    section_name = None
    page_index = None
    slide_index = None
    sheet_name = None

    for item in items:
        if item.get("line_start") is not None and (line_start is None or item["line_start"] < line_start):
            line_start = item["line_start"]
        if item.get("line_end") is not None and (line_end is None or item["line_end"] > line_end):
            line_end = item["line_end"]
        if not section_name:
            section_name = item.get("section_name")
        if page_index is None and item.get("page_index") is not None:
            page_index = item.get("page_index")
        if slide_index is None and item.get("slide_index") is not None:
            slide_index = item.get("slide_index")
        if not sheet_name and item.get("sheet_name"):
            sheet_name = item.get("sheet_name")

    source = {
        "file_name": structured_doc["file_name"],
        "file_type": structured_doc["file_type"],
        "file_id": structured_doc["file_id"],
        "highlight_text": "",
    }
    if structured_doc["file_type"] == "md" and section_name:
        source["section_name"] = section_name
    if line_start is not None:
        source["line_start"] = line_start
    if line_end is not None:
        source["line_end"] = line_end
    if page_index is not None:
        source["page_index"] = page_index
    if slide_index is not None:
        source["slide_index"] = slide_index
    if sheet_name:
        source["sheet_name"] = sheet_name
    return source


def _extract_player_stats(description: str, player_name: str) -> dict:
    stats = {}
    batting_match = re.search(
        r"(?:amassed|scored|has scored|scoring)\s+(\d[\d,]*)(?:\+)?\s+runs.*?average of ([\d.]+).*?strike rate of ([\d.]+)",
        description,
        flags=re.IGNORECASE,
    )
    if batting_match:
        stats["runs"] = _to_int(batting_match.group(1))
        stats["batting_average"] = _to_float(batting_match.group(2))
        stats["strike_rate"] = _to_float(batting_match.group(3))

    centuries_match = re.search(r"(\d+)\s+centuries?\s+and\s+(\d+)\s+fifties", description, flags=re.IGNORECASE)
    if centuries_match:
        stats["centuries"] = _to_int(centuries_match.group(1))
        stats["fifties"] = _to_int(centuries_match.group(2))
    else:
        fifties_match = re.search(r"(\d+)\s+fifties", description, flags=re.IGNORECASE)
        if fifties_match:
            stats["fifties"] = _to_int(fifties_match.group(1))

    highest_score_match = re.search(
        r"highest score(?: being| of)?\s+(\d+(?:\s+not out)?)",
        description,
        flags=re.IGNORECASE,
    )
    if highest_score_match:
        stats["highest_score"] = highest_score_match.group(1).strip()

    boundary_match = re.search(r"(\d[\d,]*)\s+fours?\s+and\s+(\d[\d,]*)\s+sixes?", description, flags=re.IGNORECASE)
    if boundary_match:
        stats["fours"] = _to_int(boundary_match.group(1))
        stats["sixes"] = _to_int(boundary_match.group(2))

    wickets_match = re.search(
        r"(?:taking|took|claimed|has taken)\s+(\d[\d,]*)\s+wickets.*?(?:average of|average)\s+([\d.]+).*?(?:economy rate of|economy)\s+([\d.]+).*?strike rate(?: of)?\s+([\d.]+)",
        description,
        flags=re.IGNORECASE,
    )
    if wickets_match:
        stats["wickets"] = _to_int(wickets_match.group(1))
        stats["bowling_average"] = _to_float(wickets_match.group(2))
        stats["economy_rate"] = _to_float(wickets_match.group(3))
        stats["bowling_strike_rate"] = _to_float(wickets_match.group(4))

    four_wicket_match = re.search(r"including\s+([A-Za-z0-9,]+)\s+four wicket haul", description, flags=re.IGNORECASE)
    if four_wicket_match:
        stats["four_wicket_hauls"] = _to_number_token(four_wicket_match.group(1))

    five_wicket_match = re.search(r"([A-Za-z0-9,]+)\s+five wicket haul", description, flags=re.IGNORECASE)
    if five_wicket_match:
        stats["five_wicket_hauls"] = _to_number_token(five_wicket_match.group(1))

    catches_match = re.search(r"(\d[\d,]*)\s+catches", description, flags=re.IGNORECASE)
    if catches_match:
        stats["catches"] = _to_int(catches_match.group(1))

    stumpings_match = re.search(r"(\d[\d,]*)\s+stumpings", description, flags=re.IGNORECASE)
    if stumpings_match:
        stats["stumpings"] = _to_int(stumpings_match.group(1))

    runouts_match = re.search(r"(\d[\d,]*)\s+runouts", description, flags=re.IGNORECASE)
    if runouts_match:
        stats["runouts"] = _to_int(runouts_match.group(1))

    if "captain" in description.lower():
        stats["is_captain"] = 1

    if "wicketkeeper" in description.lower() or "keeper" in description.lower():
        stats["is_wicketkeeper"] = 1

    stats["player_name"] = player_name
    return stats


def _extract_runs(description: str) -> int | None:
    match = _RUNS_PATTERN.search(description)
    if not match:
        return None
    try:
        return int(match.group(1).replace(",", ""))
    except ValueError:
        return None


def _clean_line(raw_line: str) -> str:
    text = (raw_line or "").strip()
    if not text:
        return ""
    if text.startswith("Got it") or text.startswith("Understood"):
        return ""
    text = re.sub(r"`{1,3}", "", text)
    text = re.sub(r"^[•*o]\s+", "", text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def _sanitize_description(text: str) -> str:
    cleaned = re.sub(r"Got it\s+[—-].*$", "", text, flags=re.IGNORECASE)
    cleaned = re.sub(r"Understood\s+[—-].*$", "", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r"\s+", " ", cleaned)
    return cleaned.strip()


def _extract_country_heading(text: str) -> str:
    cleaned = re.sub(r"[^\w\s]", " ", text or "")
    cleaned = re.sub(r"\s+", " ", cleaned).strip().lower()
    for alias, country in _COUNTRY_ALIASES.items():
        if cleaned == alias:
            return country
    for country in _COUNTRY_NAMES:
        if cleaned == country.lower():
            return country
    return ""


def _match_country_in_query(text: str) -> str:
    cleaned = re.sub(r"[^\w\s]", " ", text or "")
    cleaned = re.sub(r"\s+", " ", cleaned).strip().lower()
    for alias, country in _COUNTRY_ALIASES.items():
        if cleaned == alias or f" {alias} " in f" {cleaned} ":
            return country
    for country in _COUNTRY_NAMES:
        if cleaned == country.lower() or f" {country.lower()} " in f" {cleaned} ":
            return country
    return ""


def _is_player_heading(text: str) -> bool:
    if not text:
        return False
    if text.endswith(":"):
        return False
    if any(char.isdigit() for char in text):
        return False
    if _extract_country_heading(text):
        return False

    words = text.split()
    if len(words) < 2 or len(words) > 4:
        return False

    connector_words = {"de", "da", "di", "van", "von", "la", "del"}
    for index, word in enumerate(words):
        if word.lower() in connector_words and index not in {0, len(words) - 1}:
            continue
        if re.match(r"^[A-Z][A-Za-z'.-]*$", word):
            continue
        if re.match(r"^[A-Z]{2,}$", word):
            continue
        return False
    return True


def _ordered_countries(players: list[dict]) -> list[str]:
    countries = []
    for player in players:
        country = player["country"] or "Unknown"
        if country not in countries:
            countries.append(country)
    return countries


def _is_country_count_query(lowered_query: str) -> bool:
    return "how many" in lowered_query and "countries" in lowered_query


def _is_country_list_query(lowered_query: str) -> bool:
    if "countries" not in lowered_query:
        return False
    return "name" in lowered_query or "list" in lowered_query


def _is_stats_only_query(lowered_query: str) -> bool:
    return "only stats" in lowered_query or "full stats" in lowered_query or "stats no description" in lowered_query


def _is_definition_query(lowered_query: str) -> bool:
    return any(phrase in lowered_query for phrase in ("what is", "define", "meaning of", "explain"))


def _detect_requested_fields(query: str) -> list[str]:
    lowered_query = query.lower()
    field_map = [
        ("batting average", "batting_average"),
        ("average", "batting_average"),
        ("strike rate", "strike_rate"),
        ("runs", "runs"),
        ("centuries", "centuries"),
        ("fifties", "fifties"),
        ("highest score", "highest_score"),
        ("fours", "fours"),
        ("sixes", "sixes"),
        ("wickets", "wickets"),
        ("economy", "economy_rate"),
        ("bowling average", "bowling_average"),
        ("bowling strike rate", "bowling_strike_rate"),
        ("runouts", "runouts"),
        ("run outs", "runouts"),
        ("catches", "catches"),
        ("stumpings", "stumpings"),
        ("five wicket", "five_wicket_hauls"),
        ("fifers", "five_wicket_hauls"),
        ("four wicket", "four_wicket_hauls"),
    ]

    matched_fields = []
    for phrase, field in field_map:
        position = lowered_query.find(phrase)
        if position == -1:
            continue
        matched_fields.append((position, field))

    matched_fields.sort(key=lambda item: item[0])
    requested_fields = []
    for _, field in matched_fields:
        if field in requested_fields:
            continue
        requested_fields.append(field)
    return requested_fields


def _detect_extreme_stat_request(lowered_query: str) -> tuple[str, str] | None:
    query = lowered_query or ""
    wants_low = any(term in query for term in ("lowest", "least", "best bowling average", "best economy", "best bowling strike rate"))
    wants_high = any(term in query for term in ("highest", "most", "best", "top"))

    if not wants_low and not wants_high:
        return None

    priority = [
        ("bowling average", "bowling_average", "lowest"),
        ("economy", "economy_rate", "lowest"),
        ("bowling strike rate", "bowling_strike_rate", "lowest"),
        ("batting average", "batting_average", "highest"),
        ("average", "batting_average", "highest"),
        ("strike rate", "strike_rate", "highest"),
        ("runouts", "runouts", "highest"),
        ("run outs", "runouts", "highest"),
        ("fours", "fours", "highest"),
        ("sixes", "sixes", "highest"),
        ("runs", "runs", "highest"),
        ("wickets", "wickets", "highest"),
        ("catches", "catches", "highest"),
        ("stumpings", "stumpings", "highest"),
        ("fifers", "five_wicket_hauls", "highest"),
        ("five wicket", "five_wicket_hauls", "highest"),
        ("four wicket", "four_wicket_hauls", "highest"),
    ]
    for phrase, field, natural_direction in priority:
        if phrase in query:
            if wants_low:
                return field, "lowest"
            if wants_high:
                return field, natural_direction if "best" in query and natural_direction == "lowest" else "highest"

    if "best average" in query:
        return "batting_average", "highest"
    if "lowest average" in query:
        return "bowling_average" if "bowling" in query else "batting_average", "lowest"
    return None


def _extract_query_players(query: str, players: list[dict]) -> list[dict]:
    lowered_query = query.lower()
    matched = []
    positions = []

    for player in players:
        normalized_name = _normalize_name(player["name"])
        normalized_query = _normalize_name(query)
        if normalized_name not in normalized_query:
            continue

        position = lowered_query.find(player["name"].lower())
        if position == -1:
            position = normalized_query.find(normalized_name)
        positions.append((position, player))

    positions.sort(key=lambda item: item[0])
    for _, player in positions:
        duplicate = False
        for existing in matched:
            if existing["name"] == player["name"]:
                duplicate = True
                break
        if not duplicate:
            matched.append(player)
    return matched


def _ordered_full_stat_fields() -> list[str]:
    return [
        "runs",
        "batting_average",
        "strike_rate",
        "centuries",
        "fifties",
        "highest_score",
        "fours",
        "sixes",
        "wickets",
        "bowling_average",
        "economy_rate",
        "bowling_strike_rate",
        "four_wicket_hauls",
        "five_wicket_hauls",
        "catches",
        "stumpings",
        "runouts",
    ]


def _format_stat_value(field: str, value) -> str:
    if value is None:
        return ""
    if field == "highest_score":
        return str(value)
    if isinstance(value, int):
        return f"{value:,}"
    if isinstance(value, float):
        if value.is_integer():
            return f"{int(value)}"
        return f"{value:.2f}".rstrip("0").rstrip(".")
    return str(value)


def _field_unit(field: str) -> str:
    units = {
        "runs": "runs",
        "fours": "fours",
        "sixes": "sixes",
        "wickets": "wickets",
        "runouts": "runouts",
        "catches": "catches",
        "stumpings": "stumpings",
        "four_wicket_hauls": "4-wicket hauls",
        "five_wicket_hauls": "5-wicket hauls",
    }
    return units.get(field, "")


def _first_sentence(text: str) -> str:
    parts = re.split(r"(?<=[.!?])\s+", text.strip())
    if parts:
        return parts[0].strip()
    return text.strip()


def _normalize_name(text: str) -> str:
    cleaned = re.sub(r"[^a-z0-9]+", "", (text or "").lower())
    return cleaned


def _to_int(value: str) -> int | None:
    if not value:
        return None
    try:
        return int(value.replace(",", ""))
    except ValueError:
        return None


def _to_float(value: str) -> float | None:
    if not value:
        return None
    try:
        return float(value)
    except ValueError:
        return None


def _to_number_token(value: str) -> int | None:
    cleaned = (value or "").strip().lower()
    if not cleaned:
        return None

    digit_match = _NUMBER_TOKEN_PATTERN.search(cleaned)
    if digit_match:
        return _to_int(digit_match.group(1))

    return _WORD_NUMBERS.get(cleaned)


# =============================================================================
# CSV HANDLER  — pandas-based tabular statistics
# =============================================================================

_TABULAR_QUERY_KEYWORDS = {
    "how many rows", "row count", "number of rows", "total rows",
    "how many columns", "column count", "number of columns", "total columns",
    "what are the columns", "column names", "list columns", "list the columns",
    "average", "mean", "median", "minimum", "maximum", "min", "max",
    "total", "sum", "count", "describe", "statistics", "stats",
    "summarize", "summary", "overview", "distribution", "range",
    "missing", "null", "empty", "top values", "unique values",
}

_PDF_QUERY_KEYWORDS = {
    "how many pages", "page count", "number of pages", "total pages",
    "what are the sections", "sections", "chapters", "table of contents",
    "headings", "outline", "describe the document", "overview",
}

_DOCX_QUERY_KEYWORDS = {
    "how many pages", "how many paragraphs", "paragraph count",
    "headings", "sections", "table of contents",
    "what are the sections", "overview", "describe the document",
    "how many tables", "table count",
}

_PPTX_QUERY_KEYWORDS = {
    "how many slides", "slide count", "number of slides", "total slides",
    "slide titles", "titles", "sections", "overview", "describe the presentation",
    "summarize", "summary",
}

_IMAGE_QUERY_KEYWORDS = {
    "describe", "what does this image show", "what is in the image",
    "what is the image about", "caption", "content of image",
    "describe the image", "what can you see", "image description",
    "identify", "analyze the image",
}


def _is_tabular_query(query: str) -> bool:
    q = query.lower()
    return any(kw in q for kw in _TABULAR_QUERY_KEYWORDS)


def _is_pdf_query(query: str) -> bool:
    q = query.lower()
    return any(kw in q for kw in _PDF_QUERY_KEYWORDS)


def _is_docx_query(query: str) -> bool:
    q = query.lower()
    return any(kw in q for kw in _DOCX_QUERY_KEYWORDS)


def _is_image_query(query: str) -> bool:
    q = query.lower()
    return any(kw in q for kw in _IMAGE_QUERY_KEYWORDS)


def _is_pptx_query(query: str) -> bool:
    q = query.lower()
    return any(kw in q for kw in _PPTX_QUERY_KEYWORDS)


def _tabular_source(file_obj: File, row_start: int | None = 1, row_end: int | None = None, sheet_name: str | None = None) -> dict:
    return {
        "file_name": file_obj.original_filename or file_obj.file.name,
        "file_type": (file_obj.file_type or "").lower(),
        "file_id": file_obj.id,
        "row_start": row_start,
        "row_end": row_end,
        "sheet_name": sheet_name,
        "highlight_text": "",
    }


def _document_source(file_obj: File, page_index: int | None = 1) -> dict:
    return {
        "file_name": file_obj.original_filename or file_obj.file.name,
        "file_type": (file_obj.file_type or "").lower(),
        "file_id": file_obj.id,
        "page_index": page_index,
        "highlight_text": "",
    }


def _presentation_source(file_obj: File, slide_index: int | None = 1) -> dict:
    return {
        "file_name": file_obj.original_filename or file_obj.file.name,
        "file_type": (file_obj.file_type or "").lower(),
        "file_id": file_obj.id,
        "slide_index": slide_index,
        "highlight_text": "",
    }


def _load_dataframe_csv(file_obj: File) -> "pd.DataFrame | None":
    try:
        return pd.read_csv(file_obj.file.path, encoding="utf-8", on_bad_lines="skip")
    except Exception:
        try:
            return pd.read_csv(file_obj.file.path, encoding="latin-1", on_bad_lines="skip")
        except Exception:
            return None


def _load_dataframe_excel(file_obj: File, sheet_name: str | None = None) -> "dict[str, pd.DataFrame] | None":
    try:
        xl = pd.ExcelFile(file_obj.file.path)
        if sheet_name and sheet_name in xl.sheet_names:
            return {sheet_name: xl.parse(sheet_name)}
        return {s: xl.parse(s) for s in xl.sheet_names}
    except Exception:
        return None


def _answer_df_query(query: str, df: "pd.DataFrame", file_obj: File, label: str, row_end: int | None = None, sheet_name: str | None = None) -> dict | None:
    """
    Answer a query about a single pandas DataFrame.
    Returns a dict with 'answer' and 'sources', or None if the query is irrelevant.
    """
    q = query.lower()
    col_names = [str(c) for c in df.columns]
    col_names_lower = [c.lower() for c in col_names]
    row_count = len(df)
    col_count = len(df.columns)

    source = _tabular_source(file_obj, row_start=1, row_end=row_end or row_count, sheet_name=sheet_name)

    # ── Row / column counts ───────────────────────────────────────────────────
    if re.search(r"\bhow many rows\b|row count|number of rows|total rows", q):
        return {"answer": f"The {label} has {row_count:,} rows.", "sources": [source]}

    if re.search(r"\bhow many columns\b|column count|number of columns|total columns", q):
        return {"answer": f"The {label} has {col_count} columns.", "sources": [source]}

    if re.search(r"\bcolumn names?\b|list (?:the )?columns|what are (?:the )?columns", q):
        return {
            "answer": f"Columns in the {label} ({col_count} total): {', '.join(col_names)}.",
            "sources": [source],
        }

    # ── Summary / describe ────────────────────────────────────────────────────
    if re.search(r"\bsummar(?:y|ize|ise)\b|describe|overview|statistics\b|stats\b", q):
        lines = [f"{label} Summary", f"Rows: {row_count:,} | Columns: {col_count}",
                 f"Columns: {', '.join(col_names)}"]
        numeric_cols = df.select_dtypes(include=["number"])
        if not numeric_cols.empty:
            lines.append("\nNumeric statistics:")
            for col in numeric_cols.columns:
                series = df[col].dropna()
                if series.empty:
                    continue
                lines.append(
                    f"  {col}: min={series.min():.4g}, max={series.max():.4g}, "
                    f"mean={series.mean():.4g}, sum={series.sum():.4g}"
                )
        null_sums = df.isnull().sum()
        missing = null_sums[null_sums > 0]
        if not missing.empty:
            lines.append("\nMissing values: " + "; ".join(f"{c}: {n}" for c, n in missing.items()))
        return {"answer": "\n".join(lines), "sources": [source]}

    # ── Per-column numeric queries ────────────────────────────────────────────
    # Try to identify which column the user is asking about
    mentioned_col = None
    for i, col_lower in enumerate(col_names_lower):
        if col_lower in q:
            mentioned_col = col_names[i]
            break
        # partial match (e.g. "sales" matches "Total Sales")
        if len(col_lower) > 3 and col_lower in q:
            mentioned_col = col_names[i]
            break

    if mentioned_col:
        series = df[mentioned_col]
        if pd.api.types.is_numeric_dtype(series):
            series_clean = series.dropna()
            if re.search(r"\baverage\b|\bmean\b", q):
                return {"answer": f"Average of '{mentioned_col}': {series_clean.mean():.4g}", "sources": [source]}
            if re.search(r"\bmedian\b", q):
                return {"answer": f"Median of '{mentioned_col}': {series_clean.median():.4g}", "sources": [source]}
            if re.search(r"\bminimum\b|\bmin\b|\blowest\b|\bsmallest\b", q):
                return {"answer": f"Minimum of '{mentioned_col}': {series_clean.min():.4g}", "sources": [source]}
            if re.search(r"\bmaximum\b|\bmax\b|\bhighest\b|\blargest\b", q):
                return {"answer": f"Maximum of '{mentioned_col}': {series_clean.max():.4g}", "sources": [source]}
            if re.search(r"\btotal\b|\bsum\b", q):
                return {"answer": f"Total (sum) of '{mentioned_col}': {series_clean.sum():.4g}", "sources": [source]}
            if re.search(r"\bcount\b|\bhow many\b", q):
                non_null = series.count()
                return {"answer": f"Non-null count of '{mentioned_col}': {non_null:,}", "sources": [source]}
            if re.search(r"\bstd\b|\bstandard deviation\b", q):
                return {"answer": f"Standard deviation of '{mentioned_col}': {series_clean.std():.4g}", "sources": [source]}
            if re.search(r"\brange\b", q):
                return {
                    "answer": f"Range of '{mentioned_col}': {series_clean.min():.4g} – {series_clean.max():.4g}",
                    "sources": [source],
                }
        else:
            # Categorical
            if re.search(r"\bunique\b|distinct", q):
                unique_count = series.nunique()
                top = series.value_counts().head(5)
                top_text = ", ".join(f'"{v}": {c}' for v, c in top.items())
                return {
                    "answer": f"'{mentioned_col}' has {unique_count} unique values. Top 5: {top_text}",
                    "sources": [source],
                }
            if re.search(r"\btop\b|\bmost common\b|frequent", q):
                top = series.value_counts().head(5)
                top_text = "\n".join(f"  {v}: {c}" for v, c in top.items())
                return {
                    "answer": f"Most common values in '{mentioned_col}':\n{top_text}",
                    "sources": [source],
                }

    # ── Missing values ────────────────────────────────────────────────────────
    if re.search(r"\bmissing\b|\bnull\b|\bempty\b|NaN", q, re.IGNORECASE):
        null_sums = df.isnull().sum()
        missing = null_sums[null_sums > 0]
        if missing.empty:
            return {"answer": f"The {label} has no missing values.", "sources": [source]}
        parts = [f"{c}: {n}" for c, n in missing.items()]
        return {
            "answer": f"Missing values in the {label}:\n" + "\n".join(f"  {p}" for p in parts),
            "sources": [source],
        }

    return None


def _answer_csv_query(query: str, file_obj: File) -> dict | None:
    if not _is_tabular_query(query):
        return None
    df = _load_dataframe_csv(file_obj)
    if df is None or df.empty:
        return None
    label = file_obj.original_filename or "CSV file"
    return _answer_df_query(query, df, file_obj, label=label, row_end=len(df))


def _answer_excel_query(query: str, file_obj: File) -> dict | None:
    if not _is_tabular_query(query):
        return None
    sheets = _load_dataframe_excel(file_obj)
    if not sheets:
        return None

    # Check if user mentions a specific sheet
    q = query.lower()
    for sheet_name, df in sheets.items():
        if sheet_name.lower() in q:
            label = f'Sheet "{sheet_name}"'
            return _answer_df_query(query, df, file_obj, label=label, row_end=len(df), sheet_name=sheet_name)

    # If only one sheet, use it directly
    if len(sheets) == 1:
        sheet_name, df = next(iter(sheets.items()))
        label = file_obj.original_filename or "Excel file"
        return _answer_df_query(query, df, file_obj, label=label, row_end=len(df), sheet_name=sheet_name)

    # Multiple sheets — answer about all sheets broadly
    if re.search(r"\bhow many sheets\b|sheet count|number of sheets|sheet names", q):
        names = ", ".join(sheets.keys())
        return {
            "answer": f"The Excel file has {len(sheets)} sheet(s): {names}.",
            "sources": [_tabular_source(file_obj)],
        }

    # Try each sheet until one matches
    for sheet_name, df in sheets.items():
        label = f'Sheet "{sheet_name}"'
        result = _answer_df_query(query, df, file_obj, label=label, row_end=len(df), sheet_name=sheet_name)
        if result:
            return result

    return None


# =============================================================================
# PDF HANDLER
# =============================================================================

def _answer_pdf_query(query: str, file_obj: File) -> dict | None:
    if not _is_pdf_query(query):
        return None
    try:
        doc = fitz.open(file_obj.file.path)
        page_count = doc.page_count
        toc = doc.get_toc()
        doc.close()
    except Exception:
        return None

    q = query.lower()
    source = _document_source(file_obj, page_index=1)

    if re.search(r"\bhow many pages\b|page count|number of pages|total pages", q):
        return {"answer": f"The PDF has {page_count} page(s).", "sources": [source]}

    if re.search(r"\bsections?\b|\bchapters?\b|\boutline\b|\btable of contents\b|\bheadings?\b", q):
        if not toc:
            return {"answer": "The PDF does not have a structured table of contents or bookmarks.", "sources": [source]}
        lines = [f"The PDF has {len(toc)} section(s) / bookmarks:"]
        for level, title, page_no in toc[:25]:
            indent = "  " * (level - 1)
            lines.append(f"  {indent}{title} (page {page_no})")
        if len(toc) > 25:
            lines.append(f"  ... and {len(toc) - 25} more.")
        return {"answer": "\n".join(lines), "sources": [source]}

    if re.search(r"\bsummar(?:y|ize|ise)\b|\boverview\b|\bdescribe\b", q):
        lines = [f"PDF Document: {file_obj.original_filename or file_obj.file.name}"]
        lines.append(f"Total pages: {page_count}")
        if toc:
            lines.append(f"Sections / bookmarks: {len(toc)}")
            for level, title, page_no in toc[:10]:
                indent = "  " * (level - 1)
                lines.append(f"  {indent}{title} (page {page_no})")
        return {"answer": "\n".join(lines), "sources": [source]}

    return None


# =============================================================================
# DOCX HANDLER
# =============================================================================

def _answer_docx_query(query: str, file_obj: File) -> dict | None:
    if not _is_docx_query(query):
        return None
    try:
        doc = DocxDoc(file_obj.file.path)
        paragraphs = [p for p in doc.paragraphs if p.text.strip()]
        headings = [
            p for p in paragraphs
            if p.style and p.style.name and p.style.name.lower().startswith("heading")
        ]
        tables = doc.tables
        images = sum(1 for rel in doc.part.rels.values() if "image" in rel.reltype)
    except Exception:
        return None

    q = query.lower()
    source = _document_source(file_obj, page_index=1)

    if re.search(r"\bhow many paragraphs\b|paragraph count", q):
        return {"answer": f"The document has {len(paragraphs)} paragraph(s).", "sources": [source]}

    if re.search(r"\bhow many tables\b|table count", q):
        return {"answer": f"The document has {len(tables)} table(s).", "sources": [source]}

    if re.search(r"\bsections?\b|\bheadings?\b|\boutline\b|\btable of contents\b", q):
        if not headings:
            return {"answer": "The document does not have formal headings/sections.", "sources": [source]}
        lines = [f"The document has {len(headings)} heading(s):"]
        for p in headings[:25]:
            level = p.style.name.replace("Heading ", "H")
            lines.append(f"  [{level}] {p.text.strip()}")
        if len(headings) > 25:
            lines.append(f"  ... and {len(headings) - 25} more.")
        return {"answer": "\n".join(lines), "sources": [source]}

    if re.search(r"\bsummar(?:y|ize|ise)\b|\boverview\b|\bdescribe\b", q):
        lines = [f"Word Document: {file_obj.original_filename or file_obj.file.name}"]
        lines.append(f"Paragraphs: {len(paragraphs)}")
        lines.append(f"Tables: {len(tables)}")
        if images:
            lines.append(f"Embedded images: {images}")
        if headings:
            lines.append(f"Sections ({len(headings)} headings):")
            for p in headings[:10]:
                level = p.style.name.replace("Heading ", "H")
                lines.append(f"  [{level}] {p.text.strip()}")
        return {"answer": "\n".join(lines), "sources": [source]}

    return None


# =============================================================================
# PPTX HANDLER
# =============================================================================

def _answer_pptx_query(query: str, file_obj: File) -> dict | None:
    if not _is_pptx_query(query):
        return None

    try:
        presentation = Presentation(file_obj.file.path)
        slides = list(presentation.slides)
    except Exception:
        return None

    slide_count = len(slides)
    titles = []
    text_counts = []
    for index, slide in enumerate(slides, start=1):
        title = ""
        text_items = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text_frame.text.strip()
                if text:
                    text_items.append(text)
                    if not title:
                        title = text.splitlines()[0].strip()
        if title:
            titles.append((index, title))
        text_counts.append((index, len(text_items)))

    q = query.lower()
    source = _presentation_source(file_obj, slide_index=1)

    if re.search(r"\bhow many slides\b|slide count|number of slides|total slides", q):
        return {"answer": f"The presentation has {slide_count} slide(s).", "sources": [source]}

    if re.search(r"\bslide titles\b|\btitles\b|\bsections\b", q):
        if not titles:
            return {"answer": "The presentation slides do not have clear titles.", "sources": [source]}
        lines = [f"The presentation has {len(titles)} titled slide(s):"]
        for slide_index, title in titles[:25]:
            lines.append(f"  Slide {slide_index}: {title}")
        if len(titles) > 25:
            lines.append(f"  ... and {len(titles) - 25} more.")
        return {"answer": "\n".join(lines), "sources": [source]}

    if re.search(r"\bsummar(?:y|ize|ise)\b|\boverview\b|\bdescribe\b", q):
        lines = [f"Presentation: {file_obj.original_filename or file_obj.file.name}"]
        lines.append(f"Slides: {slide_count}")
        if titles:
            lines.append("Slide titles:")
            for slide_index, title in titles[:10]:
                lines.append(f"  Slide {slide_index}: {title}")
        if text_counts:
            busiest_slide = max(text_counts, key=lambda item: item[1])
            lines.append(f"Most text-heavy slide: Slide {busiest_slide[0]} ({busiest_slide[1]} text block(s)).")
        return {"answer": "\n".join(lines), "sources": [source]}

    return None


# =============================================================================
# IMAGE HANDLER
# =============================================================================

def _answer_image_query(query: str, file_obj: File) -> dict | None:
    """For image files, the description is already stored in Pinecone via VLM.
    We only handle metadata queries here (file size, format). Description
    queries will be handled by RAG which retrieves the VLM text from Pinecone."""
    if not _is_image_query(query):
        return None

    try:
        path = file_obj.file.path
        ext = path.rsplit(".", 1)[-1].lower()
        size_bytes = os.path.getsize(path)
        size_kb = round(size_bytes / 1024, 1)
    except Exception:
        return None

    source = _document_source(file_obj, page_index=1)
    q = query.lower()

    if re.search(r"\bformat\b|\btype\b|\bextension\b", q):
        return {
            "answer": f"The uploaded file is a {ext.upper()} image ({size_kb} KB).",
            "sources": [source],
        }

    if re.search(r"\bsize\b|\bhow large\b|\bfile size\b", q):
        return {
            "answer": f"The image file size is {size_kb} KB.",
            "sources": [source],
        }

    # For content queries, let RAG handle it (we return None so RAG kicks in)
    return None
