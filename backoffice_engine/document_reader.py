"""
document_reader.py
──────────────────
Text extraction for all supported file types.

Public API
──────────
    extract_file_text(file_path: str) -> List[Segment]

    Segment = {
        "text":         str,        # chunk-ready text for this location
        "page_index":   int | None, # PDF / DOCX / Image  (1-based)
        "slide_index":  int | None, # PPTX                (1-based)
        "sheet_name":   str | None, # XLSX
        "row_start":    int | None, # XLSX / CSV          (1-based)
        "row_end":      int | None, # XLSX / CSV          (1-based)
        "line_start":   int | None, # TXT                 (1-based)
        "line_end":     int | None, # TXT                 (1-based)
        "section_name": str | None, # Markdown
    }

VLM (Gemini)
────────────
    Used for:
      • standalone image files  (png / jpg / jpeg / webp)
      • images embedded in PDF pages
      • images embedded in DOCX  (via relationship blobs)
      • images embedded in PPTX slides  (picture shapes)
      • SVG files  (text-only prompt, no base64)

    The correct langchain_google_genai image format is:
        {"type": "image_url", "image_url": {"url": "<data-uri>"}}
    NOT the raw string directly in "image_url".
"""

import base64
import csv
import re
from itertools import groupby

import fitz                          # PyMuPDF  — PDF + image normalisation
import openpyxl                      # XLSX
import pandas as pd                  # Statistical analysis for CSV/Excel
from docx import Document as DocxDoc # DOCX
from pptx import Presentation        # PPTX
from pptx.enum.shapes import MSO_SHAPE_TYPE

from techno_chat.settings import logger
from .constants import (
    DOCX_NAMESPACE, TXT_SEGMENT_LINES, CSV_SEGMENT_ROWS, XLSX_SEGMENT_ROWS,
    MAX_RETRY_ATTEMPTS, RETRY_BACKOFF_BASE
)
from .utils import rows_to_text
from .helpers import extract_md_section_name, is_network_error, is_quota_error
from .clients import VLMClient
from .schemas import Segment
from .exceptions import (
    NetworkConnectionError, VLMQuotaExceededError,
    VLMStandaloneImageError, VLMEmbeddedImageError
)
import time

# ═══════════════════════════════════════════════════════════════════════════════
# CONSTANTS
# ═══════════════════════════════════════════════════════════════════════════════




# ═══════════════════════════════════════════════════════════════════════════════
# VLM HELPERS
# ═══════════════════════════════════════════════════════════════════════════════

def _call_vlm_embedded(image_bytes, ext, detail=""):
    """Handles VLM calls for images inside documents with retries."""
    for attempt in range(1, MAX_RETRY_ATTEMPTS + 1):
        try:
            return VLMClient().describe_image_bytes(image_bytes, ext)
        except Exception as e:
            if is_quota_error(e):
                raise VLMQuotaExceededError(internal=f"Embedded image quota. {detail}")
            if is_network_error(e):
                if attempt < MAX_RETRY_ATTEMPTS:
                    wait = RETRY_BACKOFF_BASE * (2 ** (attempt - 1))
                    logger.warning("Network retry (embedded image) | attempt=%s/%s | wait=%s | %s",
                                   attempt, MAX_RETRY_ATTEMPTS, wait, detail)
                    time.sleep(wait)
                    continue
                raise NetworkConnectionError(internal=f"Network on embedded image after retries. {detail}")
            raise VLMEmbeddedImageError(internal=f"VLM failed on embedded image. {detail} raw={e}")


# ═══════════════════════════════════════════════════════════════════════════════
# CONSTANTS
# ═══════════════════════════════════════════════════════════════════════════════




# ═══════════════════════════════════════════════════════════════════════════════
# SEGMENT FACTORY
# ═══════════════════════════════════════════════════════════════════════════════

def _seg(text: str, **loc) -> Segment:
    """
    Build one Segment dict.
    All eight location fields default to None; pass only what is known.

    Example:
        _seg("hello world", page_index=3)
        _seg("slide text",  slide_index=2)
        _seg("row data",    sheet_name="Revenue", row_start=1, row_end=50)
    """
    return {
        "text":         text,
        "page_index":   loc.get("page_index"),
        "slide_index":  loc.get("slide_index"),
        "sheet_name":   loc.get("sheet_name"),
        "row_start":    loc.get("row_start"),
        "row_end":      loc.get("row_end"),
        "line_start":   loc.get("line_start"),
        "line_end":     loc.get("line_end"),
        "section_name": loc.get("section_name"),
    }


# ═══════════════════════════════════════════════════════════════════════════════
# VLM  —  Gemini image understanding
# ═══════════════════════════════════════════════════════════════════════════════







# ═══════════════════════════════════════════════════════════════════════════════
# TABLE HELPERS  (reused by PDF / DOCX / PPTX)
# ═══════════════════════════════════════════════════════════════════════════════




# ═══════════════════════════════════════════════════════════════════════════════
# PDF
# ═══════════════════════════════════════════════════════════════════════════════

def _pdf_page_parts(page, page_number: int) -> list[str]:
    """
    Extract all content from one PDF page in reading order.

    Text   -> get_text blocks sorted top-to-bottom.
    Images -> page.get_images() + doc.extract_image(xref) for real bytes.
              block["image"] from get_text("dict") is empty for vector/referenced
              images — xref extraction is the only reliable method.
    Tables -> PyMuPDF table finder.
    """
    doc   = page.parent              # ← fitz.Document — needed for extract_image
    parts = []

    # ── 1. Text in reading order ──────────────────────────────────────────────
    blocks = sorted(
        page.get_text("dict")["blocks"],
        key=lambda b: (round(b["bbox"][1], 1), b["bbox"][0]),
    )
    for block in blocks:
        if block["type"] == 0:                         # text block only
            for line in block["lines"]:
                text = "".join(s["text"] for s in line.get("spans", [])).strip()
                if text:
                    parts.append(text)
    # type==1 image blocks intentionally skipped — handled below via xref

    # ── 2. Images via xref ────────────────────────────────────────────────────

    seen_xrefs    = set()            # ← outside the loop, resets once per page
    image_counter = 0                # ← outside the loop, counts per page

    for img_info in page.get_images(full=True):
        xref = img_info[0]
        if xref in seen_xrefs:
            continue
        seen_xrefs.add(xref)
        image_counter += 1

        try:
            extracted = doc.extract_image(xref)  # ← uses doc defined above
            if not extracted:
                logger.warning("PDF image xref=%s returned nothing  page=%s",
                               xref, page_number)
                continue

            image_bytes = extracted["image"]
            ext         = extracted.get("ext", "png")

            if not image_bytes:
                logger.warning("PDF image xref=%s empty bytes  page=%s",
                                   xref, page_number)
                continue

            ocr = _call_vlm_embedded(image_bytes, ext, detail=f"page={page_number} xref={xref}")
            if ocr:
                parts.append(f"[Image {image_counter}: {ocr}]")

        except VLMEmbeddedImageError as exc:
            logger.warning("PDF image OCR failed (skipped) | %s", exc.internal_note)
        except (NetworkConnectionError, VLMQuotaExceededError):
            raise

    # ── 3. Tables ─────────────────────────────────────────────────────────────
    try:
        for table in page.find_tables():
            data = table.extract()
            if data:
                parts.append(rows_to_text(data))
    except Exception as exc:
        logger.warning("PDF table extraction failed  page=%s  %s", page_number, exc)

    return parts


def pdf_file_text(file_path: str) -> list:
    doc      = fitz.open(file_path)
    segments = []

    for page_number, page in enumerate(doc, start=1):
        parts = _pdf_page_parts(page, page_number)
        text  = "\n".join(parts).strip()
        if text:
            segments.append(_seg(text, page_index=page_number))

    doc.close()
    logger.info("pdf_file_text | pages=%s", len(segments))
    return segments


# ═══════════════════════════════════════════════════════════════════════════════
# DOCX
# ═══════════════════════════════════════════════════════════════════════════════

def _docx_count_page_breaks(para_elem) -> int:
    """
    Count page-break signals inside one paragraph XML element.
    Handles both rendered breaks and explicit breaks.
    """
    W_PB = f"{{{DOCX_NAMESPACE}}}lastRenderedPageBreak"   # <w:lastRenderedPageBreak/>
    W_BR = f"{{{DOCX_NAMESPACE}}}br"                      # <w:br w:type="page"/>
    count = 0
    for elem in para_elem.iter():
        if elem.tag == W_PB:
            count += 1
        elif elem.tag == W_BR and elem.get(f"{{{DOCX_NAMESPACE}}}type") == "page":
            count += 1
    return count


def _docx_elem_text(elem) -> str:
    """Collect all <w:t> run text from any XML element."""
    return "".join(
        node.text or ""
        for node in elem.iter()
        if node.tag == f"{{{DOCX_NAMESPACE}}}t"
    ).strip()


def docx_file_text(file_path: str) -> list:
    """
    Page detection:
      Walk body elements in document order.
      Increment page counter on every explicit / rendered page break.
      Tables and paragraphs are tagged with the current page number.

      Images cannot be reliably placed on a page without rendering,
      so they are appended to page 1 with clear labels.
    """
    doc           = DocxDoc(file_path)
    body          = doc.element.body
    lines_by_page = []          # list of (text, page_number)
    current_page  = 1
    image_counter = 0

    for child in body:
        local = child.tag.split("}")[-1] if "}" in child.tag else child.tag

        if local == "p":                               # paragraph
            current_page += _docx_count_page_breaks(child)
            text = _docx_elem_text(child)
            if text:
                lines_by_page.append((text, current_page))

        elif local == "tbl":                           # table
            rows = []
            for tr in child.iter(f"{{{DOCX_NAMESPACE}}}tr"):
                cells = [_docx_elem_text(tc) for tc in tr.iter(f"{{{DOCX_NAMESPACE}}}tc")]
                if any(cells):
                    rows.append(cells)
            if rows:
                lines_by_page.append((rows_to_text(rows), current_page))

    # Inline images (via part relationships — position unknown, put on page 1)
    for rel in doc.part.rels.values():
        if "image" not in rel.reltype:
            continue
        image_counter += 1
        try:
            blob = rel.target_part.blob
            ext  = rel.target_part.content_type.split("/")[-1]
            ocr  = _call_vlm_embedded(blob, ext, detail=f"docx_image={image_counter}")
            lines_by_page.append((f"[Image {image_counter}: {ocr}]", 1))
        except VLMEmbeddedImageError as exc:
            logger.warning("DOCX image OCR failed (skipped) | %s", exc.internal_note)
        except (NetworkConnectionError, VLMQuotaExceededError):
            raise

    if not lines_by_page:
        logger.info("docx_file_text | no content")
        return []

    # Group consecutive lines that share the same page number
    segments = []
    for page_num, group in groupby(lines_by_page, key=lambda x: x[1]):
        text = "\n".join(line for line, _ in group).strip()
        if text:
            segments.append(_seg(text, page_index=page_num))

    logger.info("docx_file_text | pages=%s", len(segments))
    return segments


# ═══════════════════════════════════════════════════════════════════════════════
# PPTX
# ═══════════════════════════════════════════════════════════════════════════════

def pptx_file_text(file_path: str) -> list:
    """
    One segment per slide.
    Text, tables, and images are all extracted inline for that slide.
    """
    prs           = Presentation(file_path)
    segments      = []
    image_counter = 0

    for slide_no, slide in enumerate(prs.slides, start=1):
        parts = [f"[Slide {slide_no}]"]

        for shape in slide.shapes:
            # ── Text ──
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)

            # ── Table ──
            elif shape.has_table:
                rows = [
                    [cell.text.strip() for cell in row.cells]
                    for row in shape.table.rows
                ]
                parts.append(rows_to_text(rows))

            # ── Picture ──
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_counter += 1
                try:
                    ocr = _call_vlm_embedded(shape.image.blob, shape.image.ext, detail=f"slide={slide_no} img={image_counter}")
                    parts.append(f"[Image {image_counter}: {ocr}]")
                except VLMEmbeddedImageError as exc:
                    logger.warning("PPTX image OCR failed (skipped) | %s", exc.internal_note)
                except (NetworkConnectionError, VLMQuotaExceededError):
                    raise

        text = "\n".join(parts).strip()
        if text:
            segments.append(_seg(text, slide_index=slide_no))

    logger.info("pptx_file_text | slides=%s", len(segments))
    return segments


# ═══════════════════════════════════════════════════════════════════════════════
# XLSX
# ═══════════════════════════════════════════════════════════════════════════════

def xlsx_file_text(file_path: str) -> list:
    """
    One or more segments per sheet (up to XLSX_SEGMENT_ROWS data rows each).
    sheet_name / row_start / row_end are stored per segment.
    """
    wb = openpyxl.load_workbook(file_path, data_only=True)
    segments = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = [list(row) for row in ws.iter_rows(values_only=True)]
        non_empty_rows = [row for row in rows if any(cell is not None and str(cell).strip() for cell in row)]
        if not non_empty_rows:
            continue

        header_row = ["" if cell is None else str(cell).strip() for cell in non_empty_rows[0]]
        has_header = any(header_row)
        data_rows = non_empty_rows[1:] if has_header and len(non_empty_rows) > 1 else non_empty_rows
        row_offset = 2 if has_header and len(non_empty_rows) > 1 else 1

        for i in range(0, len(data_rows), XLSX_SEGMENT_ROWS):
            batch = data_rows[i : i + XLSX_SEGMENT_ROWS]
            rendered_rows = []
            for row_index, row in enumerate(batch, start=i + row_offset):
                cells = []
                for column_index, cell in enumerate(row):
                    value = "" if cell is None else str(cell).strip()
                    if not value:
                        continue
                    header = header_row[column_index] if column_index < len(header_row) else ""
                    cells.append(f"{header}: {value}" if header else value)
                if cells:
                    rendered_rows.append(f"Row {row_index}: " + " | ".join(cells))

            if not rendered_rows:
                continue

            prefix = [f"[Sheet: {sheet_name}]"]
            if has_header:
                prefix.append("Columns: " + ", ".join(item for item in header_row if item))
            text = "\n".join(prefix + rendered_rows)
            segments.append(_seg(text.strip(), sheet_name=sheet_name, row_start=i + row_offset, row_end=i + row_offset + len(batch) - 1))

    logger.info("xlsx_file_text | segments=%s", len(segments))
    return segments


# ═══════════════════════════════════════════════════════════════════════════════
# CSV
# ═══════════════════════════════════════════════════════════════════════════════

def csv_file_text(file_path: str) -> list:
    """
    Groups rows into CSV_SEGMENT_ROWS-row segments.
    Row numbers are 1-based (row 1 = first row, usually the header).
    """
    all_rows = []
    with open(file_path, "r", encoding="utf-8", errors="replace", newline="") as fh:
        for row in csv.reader(fh):
            all_rows.append([str(cell).strip() for cell in row])

    segments = []
    if not all_rows:
        return segments

    header_row = all_rows[0]
    has_header = any(header_row)
    data_rows = all_rows[1:] if has_header and len(all_rows) > 1 else all_rows
    row_offset = 2 if has_header and len(all_rows) > 1 else 1

    for i in range(0, len(data_rows), CSV_SEGMENT_ROWS):
        batch = data_rows[i : i + CSV_SEGMENT_ROWS]
        rendered_rows = []
        for row_index, row in enumerate(batch, start=i + row_offset):
            cells = []
            for column_index, value in enumerate(row):
                if not value:
                    continue
                header = header_row[column_index] if column_index < len(header_row) else ""
                cells.append(f"{header}: {value}" if header else value)
            if cells:
                rendered_rows.append(f"Row {row_index}: " + " | ".join(cells))

        if not rendered_rows:
            continue

        prefix = []
        if has_header:
            prefix.append("Columns: " + ", ".join(item for item in header_row if item))
        text = "\n".join(prefix + rendered_rows).strip()
        if text:
            segments.append(_seg(text, row_start=i + row_offset, row_end=i + row_offset + len(batch) - 1))

    logger.info("csv_file_text | segments=%s", len(segments))
    return segments


# ═══════════════════════════════════════════════════════════════════════════════
# TXT
# ═══════════════════════════════════════════════════════════════════════════════

def txt_file_text(file_path: str) -> list:
    """
    Segments of TXT_SEGMENT_LINES lines each.
    line_start / line_end record the 1-based range for the stored block.
    Frontend shows a 50-line display window (line_start to line_start + 49).
    """
    with open(file_path, "r", encoding="utf-8", errors="replace") as fh:
        lines = fh.readlines()

    segments = []
    for i in range(0, len(lines), TXT_SEGMENT_LINES):
        batch = lines[i : i + TXT_SEGMENT_LINES]
        text  = "".join(batch).strip()
        if text:
            segments.append(_seg(
                text,
                line_start=i + 1,
                line_end=i + len(batch),
            ))

    logger.info("txt_file_text | segments=%s", len(segments))
    return segments


# ═══════════════════════════════════════════════════════════════════════════════
# MARKDOWN
# ═══════════════════════════════════════════════════════════════════════════════




def md_file_text(file_path: str) -> list:
    """
    Splits on any heading level (# through ######).
    Each section's section_name is the heading text.
    Image markdown is replaced with a reference tag.
    """
    with open(file_path, "r", encoding="utf-8", errors="replace") as fh:
        content = fh.read()

    # Replace embedded image syntax with a readable reference
    content = re.sub(
        r"!\[([^\]]*)\]\(([^)]+)\)",
        lambda m: f"[Image Ref: alt='{m.group(1).strip()}' url='{m.group(2).strip()}']",
        content,
    )

    raw_sections = re.split(r"(?m)^(?=#{1,6} )", content)
    segments     = []

    for i, section in enumerate(raw_sections, start=1):
        text = section.strip()
        if text:
            segments.append(_seg(
                text,
                section_name=extract_md_section_name(text, fallback=f"Section {i}"),
            ))

    logger.info("md_file_text | sections=%s", len(segments))
    return segments


# ═══════════════════════════════════════════════════════════════════════════════
# IMAGE FILES  (standalone upload)
# ═══════════════════════════════════════════════════════════════════════════════

def image_file_text(file_path: str) -> list:
    """
    Standalone image file (PNG / JPG / JPEG / WEBP / SVG).
    page_index is always 1 — a single image is a single page.
    """
    ext = file_path.rsplit(".", 1)[-1].lower()
    
    for attempt in range(1, MAX_RETRY_ATTEMPTS + 1):
        try:
            if ext == "svg":
                text = VLMClient().describe_svg_file(file_path)
            else:
                text = VLMClient().describe_image_file(file_path)
            
            logger.info("image_file_text | chars=%s", len(text))
            return [_seg(text, page_index=1)] if text else []
            
        except Exception as e:
            if is_network_error(e) and attempt < MAX_RETRY_ATTEMPTS:
                wait = RETRY_BACKOFF_BASE * (2 ** (attempt - 1))
                logger.warning(
                    "Network retry (standalone image) | attempt=%s/%s | wait=%s | file=%s",
                    attempt,
                    MAX_RETRY_ATTEMPTS,
                    wait,
                    file_path,
                )
                time.sleep(wait)
                continue

            logger.warning("image_file_text | VLM unavailable, using metadata fallback | file=%s raw=%s", file_path, e)
            fallback_text = (
                f"Image file uploaded successfully. Format: {ext.upper()}. "
                "Detailed visual analysis is currently unavailable, but the file is available in the document library."
            )
            return [_seg(fallback_text, page_index=1)]

    return []


# ═══════════════════════════════════════════════════════════════════════════════
# STATISTICAL SUMMARY SEGMENT GENERATORS
# These produce a human-readable stats block prepended to every file's segments
# so queries like "how many rows", "what columns", "page count" are answered.
# ═══════════════════════════════════════════════════════════════════════════════

def _csv_stats_segment(file_path: str) -> Segment | None:
    """
    Build a stats-summary segment for a CSV file using pandas.
    Returns one Segment with row count, column names, and numeric stats.
    """
    try:
        df = pd.read_csv(file_path, encoding="utf-8", on_bad_lines="skip")
        return _dataframe_stats_segment(df, label="CSV File")
    except Exception as exc:
        logger.warning("_csv_stats_segment | failed: %s", exc)
        return None


def _excel_stats_segment(file_path: str) -> Segment | None:
    """
    Build a stats-summary segment for an Excel file using pandas.
    Combines stats from all sheets into one overview segment.
    """
    try:
        xl = pd.ExcelFile(file_path)
        sheet_summaries = []
        for sheet in xl.sheet_names:
            df = xl.parse(sheet)
            summary = _dataframe_summary_text(df, label=f'Sheet "{sheet}"')
            sheet_summaries.append(summary)
        full_text = f"Excel File Statistics\nSheets: {', '.join(xl.sheet_names)}\n\n" + "\n\n".join(sheet_summaries)
        return _seg(full_text.strip(), sheet_name=xl.sheet_names[0] if xl.sheet_names else None, row_start=1)
    except Exception as exc:
        logger.warning("_excel_stats_segment | failed: %s", exc)
        return None


def _dataframe_stats_segment(df: "pd.DataFrame", label: str) -> Segment:
    """Convert a DataFrame into a Segment with full statistical summary."""
    text = _dataframe_summary_text(df, label=label)
    return _seg(text.strip(), row_start=1, row_end=len(df))


def _dataframe_summary_text(df: "pd.DataFrame", label: str) -> str:
    """
    Generate a rich plain-text statistical summary for a DataFrame,
    formatted so the LLM can answer questions like:
      - how many rows / columns?
      - what are the column names?
      - what is the average/min/max of column X?
      - what is the total of column X?
      - which column has the most NaN values?
    """
    lines = []
    lines.append(f"{label} Statistical Summary")
    lines.append(f"Total rows: {len(df)}")
    lines.append(f"Total columns: {len(df.columns)}")
    lines.append(f"Column names: {', '.join(str(c) for c in df.columns)}")

    # Null / missing values per column
    null_counts = df.isnull().sum()
    non_null = null_counts[null_counts > 0]
    if not non_null.empty:
        null_parts = [f"{col}: {cnt} missing" for col, cnt in non_null.items()]
        lines.append(f"Columns with missing values: {'; '.join(null_parts)}")
    else:
        lines.append("No missing values.")

    # Numeric columns — describe
    numeric_cols = df.select_dtypes(include=["number"])
    if not numeric_cols.empty:
        lines.append("\nNumeric Column Statistics:")
        for col in numeric_cols.columns:
            series = df[col].dropna()
            if series.empty:
                continue
            lines.append(
                f"  {col}: "
                f"min={series.min():.4g}, "
                f"max={series.max():.4g}, "
                f"mean={series.mean():.4g}, "
                f"median={series.median():.4g}, "
                f"sum={series.sum():.4g}, "
                f"std={series.std():.4g}"
            )

    # Categorical columns — value counts (top 5)
    cat_cols = df.select_dtypes(exclude=["number"])
    if not cat_cols.empty:
        lines.append("\nCategorical Columns (top 5 values each):")
        for col in cat_cols.columns[:10]:  # limit to 10 categorical cols
            vc = df[col].value_counts().head(5)
            top_vals = ", ".join(f'"{v}": {c}' for v, c in vc.items())
            lines.append(f"  {col}: {top_vals}")

    return "\n".join(lines)


def _pdf_stats_segment(file_path: str) -> Segment | None:
    """
    Build a stats-summary segment for a PDF: page count + any bookmarks/outline.
    """
    try:
        doc = fitz.open(file_path)
        page_count = doc.page_count

        lines = ["PDF Document Statistics"]
        lines.append(f"Total pages: {page_count}")

        # Outline / table of contents
        toc = doc.get_toc()
        if toc:
            lines.append(f"Sections / bookmarks: {len(toc)}")
            for level, title, page_no in toc[:20]:  # show first 20
                indent = "  " * (level - 1)
                lines.append(f"  {indent}{title} (page {page_no})")
            if len(toc) > 20:
                lines.append(f"  ... and {len(toc) - 20} more sections")
        else:
            lines.append("No bookmark/outline structure found.")

        # Sample first page word count
        if page_count > 0:
            first_page_text = doc[0].get_text("text")
            word_count = len(first_page_text.split())
            lines.append(f"First page word count (approx): {word_count}")

        doc.close()
        return _seg("\n".join(lines).strip(), page_index=1)
    except Exception as exc:
        logger.warning("_pdf_stats_segment | failed: %s", exc)
        return None


def _docx_stats_segment(file_path: str) -> Segment | None:
    """
    Build a stats-summary segment for a DOCX: paragraph count, heading structure.
    """
    try:
        doc = DocxDoc(file_path)
        paragraphs = [p for p in doc.paragraphs if p.text.strip()]
        headings = [
            p for p in paragraphs
            if p.style and p.style.name and p.style.name.lower().startswith("heading")
        ]
        tables = doc.tables

        lines = ["Word Document Statistics"]
        lines.append(f"Total paragraphs: {len(paragraphs)}")
        lines.append(f"Total tables: {len(tables)}")

        if headings:
            lines.append(f"Headings / sections ({len(headings)} total):")
            for heading in headings[:20]:
                level = heading.style.name.replace("Heading ", "").strip()
                lines.append(f"  [H{level}] {heading.text.strip()}")
            if len(headings) > 20:
                lines.append(f"  ... and {len(headings) - 20} more headings")
        else:
            lines.append("No formal heading structure found.")

        # Inline images count
        image_count = sum(
            1 for rel in doc.part.rels.values()
            if "image" in rel.reltype
        )
        if image_count:
            lines.append(f"Embedded images: {image_count}")

        return _seg("\n".join(lines).strip(), page_index=1)
    except Exception as exc:
        logger.warning("_docx_stats_segment | failed: %s", exc)
        return None


def _image_stats_segment(file_path: str, vlm_description: str) -> Segment | None:
    """
    Build a stats-summary segment for an image file using its VLM description.
    """
    try:
        ext = file_path.rsplit(".", 1)[-1].lower()
        import os
        size_bytes = os.path.getsize(file_path)
        size_kb = round(size_bytes / 1024, 1)

        lines = ["Image File Summary"]
        lines.append(f"Format: {ext.upper()}")
        lines.append(f"File size: {size_kb} KB")
        if vlm_description:
            lines.append(f"\nContents:\n{vlm_description}")

        return _seg("\n".join(lines).strip(), page_index=1)
    except Exception as exc:
        logger.warning("_image_stats_segment | failed: %s", exc)
        return None


def _txt_stats_segment(file_path: str) -> Segment | None:
    """
    Build a stats-summary segment for a TXT file: line count, word count, etc.
    """
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as fh:
            lines_raw = fh.readlines()
        line_count = len(lines_raw)
        word_count = sum(len(l.split()) for l in lines_raw)
        char_count = sum(len(l) for l in lines_raw)
        non_empty = sum(1 for l in lines_raw if l.strip())

        text = (
            f"Text File Statistics\n"
            f"Total lines: {line_count}\n"
            f"Non-empty lines: {non_empty}\n"
            f"Total words (approx): {word_count}\n"
            f"Total characters: {char_count}"
        )
        return _seg(text.strip(), line_start=1, line_end=line_count)
    except Exception as exc:
        logger.warning("_txt_stats_segment | failed: %s", exc)
        return None


def _md_stats_segment(file_path: str) -> Segment | None:
    """
    Build a stats-summary segment for a Markdown file:
    heading count, section count, line count, word count.
    """
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as fh:
            content = fh.read()
        lines_raw = content.splitlines()
        line_count = len(lines_raw)
        word_count = len(content.split())
        headings = [l.strip() for l in lines_raw if re.match(r"^#{1,6}\s", l)]
        sections = len(headings)

        text_lines = [
            "Markdown Document Statistics",
            f"Total lines: {line_count}",
            f"Total words (approx): {word_count}",
            f"Headings / sections: {sections}",
        ]
        if headings:
            text_lines.append("Sections:")
            for heading in headings[:20]:
                text_lines.append(f"  {heading}")
            if len(headings) > 20:
                text_lines.append(f"  ... and {len(headings) - 20} more")

        return _seg("\n".join(text_lines).strip(), section_name="Document Statistics")
    except Exception as exc:
        logger.warning("_md_stats_segment | failed: %s", exc)
        return None


def _pptx_stats_segment(file_path: str) -> Segment | None:
    """
    Build a stats-summary segment for a PPTX: slide count, slide titles.
    """
    try:
        prs = Presentation(file_path)
        slide_count = len(prs.slides)
        titles = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and hasattr(shape, "placeholder_format"):
                    if shape.placeholder_format and shape.placeholder_format.idx == 0:
                        title = shape.text_frame.text.strip()
                        if title:
                            titles.append(title)
                        break

        lines = ["PowerPoint Presentation Statistics"]
        lines.append(f"Total slides: {slide_count}")
        if titles:
            lines.append("Slide titles:")
            for i, title in enumerate(titles[:30], start=1):
                lines.append(f"  Slide {i}: {title}")
        return _seg("\n".join(lines).strip(), slide_index=1)
    except Exception as exc:
        logger.warning("_pptx_stats_segment | failed: %s", exc)
        return None


def _xlsx_stats_segment(file_path: str) -> Segment | None:
    """Delegate to _excel_stats_segment."""
    return _excel_stats_segment(file_path)


# ═══════════════════════════════════════════════════════════════════════════════
# ROUTER
# ═══════════════════════════════════════════════════════════════════════════════

_EXTRACTOR_MAP: dict = {
    "pdf":  pdf_file_text,
    "docx": docx_file_text,
    "doc":  docx_file_text,
    "txt":  txt_file_text,
    "md":   md_file_text,
    "csv":  csv_file_text,
    "png":  image_file_text,
    "jpg":  image_file_text,
    "jpeg": image_file_text,
    "webp": image_file_text,
    "svg":  image_file_text,
    "pptx": pptx_file_text,
    "xlsx": xlsx_file_text,
}

# Maps extension → stats segment generator
_STATS_SEGMENT_MAP: dict = {
    "csv":  _csv_stats_segment,
    "xlsx": _xlsx_stats_segment,
    "pdf":  _pdf_stats_segment,
    "docx": _docx_stats_segment,
    "doc":  _docx_stats_segment,
    "txt":  _txt_stats_segment,
    "md":   _md_stats_segment,
    "pptx": _pptx_stats_segment,
}


def extract_file_text(file_path: str) -> list:
    """
    Universal entry point.

    Detects file extension → calls the correct extractor →
    prepends a statistical-summary segment (if available) →
    returns List[Segment].

    Each Segment dict:
        {
            "text":         str,
            "page_index":   int | None,
            "slide_index":  int | None,
            "sheet_name":   str | None,
            "row_start":    int | None,
            "row_end":      int | None,
            "line_start":   int | None,
            "line_end":     int | None,
            "section_name": str | None,
        }

    Raises ValueError for unsupported extensions.
    """
    ext       = file_path.rsplit(".", 1)[-1].lower()
    extractor = _EXTRACTOR_MAP.get(ext)
    if extractor is None:
        raise ValueError(f"Unsupported file type: '.{ext}'")

    segments = extractor(file_path)

    # ── Prepend stats-summary segment (standalone images handled specially) ──
    if ext in ("png", "jpg", "jpeg", "webp", "svg"):
        # For standalone images we already have the VLM description in segments[0]
        vlm_text = segments[0]["text"] if segments else ""
        stats_seg = _image_stats_segment(file_path, vlm_text)
    else:
        stats_fn = _STATS_SEGMENT_MAP.get(ext)
        stats_seg = stats_fn(file_path) if stats_fn else None

    if stats_seg and stats_seg.get("text"):
        # Deduplicate: don't add if text is already the first segment
        if not segments or segments[0].get("text", "").strip() != stats_seg["text"].strip():
            segments = [stats_seg] + list(segments)
            logger.info("extract_file_text | stats segment prepended | ext=%s", ext)

    return segments
