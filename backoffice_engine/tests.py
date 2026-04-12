import json
import shutil
import tempfile
from pathlib import Path
from types import SimpleNamespace
from unittest.mock import MagicMock, patch

import fitz
import openpyxl
from django.conf import settings
from django.contrib.admin.models import CHANGE, LogEntry
from django.contrib.messages import get_messages
from django.core.files.uploadedfile import SimpleUploadedFile
from django.http import HttpResponse
from django.test import SimpleTestCase, TestCase, override_settings
from django.urls import reverse
from docx import Document as DocxDocument
from docx.shared import Inches
from PIL import Image
from pptx import Presentation
from pptx.util import Inches as PptxInches

from backoffice_engine import ai_assistant_service, chat_service, document_reader, image_generation_service, web_search_service
from backoffice_engine.admin_auth import create_admin
from backoffice_engine.choices import FileProcessingStatus, FileType, SessionType
from backoffice_engine.conversation_state_service import get_last_active_chat_session_id
from backoffice_engine.document_reader import extract_file_text
from backoffice_engine.exceptions import ChatResponseError, NoTextExtractedError, WebSearchError
from backoffice_engine.ingestion_service import embed_file_and_upsert
from backoffice_engine.models import AdminProfile, AdminUser, ChatMessage, ChatSession, File, User, UserProfile


class FakeLLMResponse:
    def __init__(self, content: str):
        self.content = content


class FakeLLM:
    def __init__(self, responses):
        self.responses = list(responses)
        self.calls = []
        self.last_response = self.responses[-1] if self.responses else ""

    def invoke(self, payload):
        self.calls.append(payload)
        response = self.responses.pop(0) if self.responses else self.last_response
        return FakeLLMResponse(response)


class FakePromptPipe:
    def __or__(self, llm):
        return llm


def fake_render_factory(store, content="ok"):
    def _fake_render(_request, template_name, context=None, *args, **kwargs):
        store["template_name"] = template_name
        store["context"] = context or {}
        return HttpResponse(content)

    return _fake_render


class FakeDenseItem:
    def __init__(self, values):
        self.values = values


class FakeDenseEmbeddings:
    def __init__(self, count):
        self.data = [FakeDenseItem([0.1, 0.2, 0.3]) for _ in range(count)]


class FakeSparseItem:
    def __init__(self, indices=None, values=None):
        self._indices = indices or [1, 2, 3]
        self._values = values or [0.2, 0.4, 0.6]

    def to_dict(self):
        return {
            "sparse_indices": self._indices,
            "sparse_values": self._values,
        }


class FakeSparseEmbeddings:
    def __init__(self, count):
        self.data = [FakeSparseItem() for _ in range(count)]


class BaseTechnoChatTestCase(TestCase):
    def setUp(self):
        super().setUp()
        self.temp_media_dir = tempfile.mkdtemp(prefix="technochat-tests-")
        override = override_settings(
            MEDIA_ROOT=self.temp_media_dir,
            ALLOWED_HOSTS=["testserver", "localhost"],
        )
        override.enable()
        self.addCleanup(override.disable)
        self.addCleanup(lambda: shutil.rmtree(self.temp_media_dir, ignore_errors=True))

    def create_contributor(
        self,
        email="user@technostacks.com",
        password="StrongPass1!",
        profile_complete=True,
        first_name="User",
        surname="Tester",
        username="user.tester",
        position="QA Engineer",
        team="QA",
    ):
        user = User.objects.create(
            email=email,
            password=password,
            profile_completed=profile_complete,
        )
        profile = UserProfile.objects.create(
            user=user,
            first_name=first_name,
            surname=surname,
            username=username,
            position_at_technostacks=position,
            team=team,
            is_profile_complete=profile_complete,
        )
        return user, profile

    def create_admin_user(
        self,
        email="admin@technostacks.com",
        password="StrongPass1!",
        profile_complete=True,
        is_staff=True,
        first_name="Admin",
        surname="Tester",
        username="admin.tester",
        team="Core",
    ):
        admin = AdminUser(
            email=email,
            is_staff=is_staff,
            is_superuser=is_staff,
            profile_completed=profile_complete,
        )
        admin.set_password(password)
        admin.save()
        profile = AdminProfile.objects.create(
            admin=admin,
            first_name=first_name,
            surname=surname,
            username=username,
            position_at_technostacks="Administrator",
            team=team,
            is_profile_complete=profile_complete,
        )
        return admin, profile

    def set_contributor_session(self, user):
        session = self.client.session
        session["user_id"] = user.id
        session.save()

    def set_admin_session(self, admin):
        self.client.force_login(admin)
        session = self.client.session
        session["tc_admin_id"] = admin.id
        session["role"] = "admin"
        session.save()

    def image_file(self, name="image.png", format="PNG", color=(220, 130, 40)):
        path = Path(self.temp_media_dir) / name
        image = Image.new("RGB", (40, 40), color)
        image.save(path, format=format)
        return path

    def pdf_file(self, name="sample.pdf", pages=None, include_image=False):
        path = Path(self.temp_media_dir) / name
        doc = fitz.open()
        pages = pages or ["PDF page one", "PDF page two"]
        image_path = self.image_file("embedded.png")
        for index, text in enumerate(pages, start=1):
            page = doc.new_page()
            page.insert_text((72, 72), text)
            if include_image and index == 1:
                page.insert_image(fitz.Rect(72, 120, 180, 220), filename=str(image_path))
        doc.save(path)
        doc.close()
        return path

    def docx_file(self, name="sample.docx", with_table=False, with_image=False):
        path = Path(self.temp_media_dir) / name
        doc = DocxDocument()
        doc.add_heading("Project Notes", level=1)
        doc.add_paragraph("Document body for DOCX extraction tests.")
        if with_table:
            table = doc.add_table(rows=2, cols=2)
            table.rows[0].cells[0].text = "Name"
            table.rows[0].cells[1].text = "Role"
            table.rows[1].cells[0].text = "Raj"
            table.rows[1].cells[1].text = "QA"
        if with_image:
            doc.add_picture(str(self.image_file("docx-embedded.png")), width=Inches(1))
        doc.save(path)
        return path

    def pptx_file(self, name="sample.pptx", with_table=False, with_image=False):
        path = Path(self.temp_media_dir) / name
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        title_box = slide.shapes.add_textbox(PptxInches(0.6), PptxInches(0.4), PptxInches(4), PptxInches(0.7))
        title_box.text_frame.text = "Sprint Review"
        body_box = slide.shapes.add_textbox(PptxInches(0.6), PptxInches(1.4), PptxInches(5), PptxInches(1.2))
        body_box.text_frame.text = "Presentation body text"
        if with_table:
            table = slide.shapes.add_table(2, 2, PptxInches(0.6), PptxInches(2.8), PptxInches(4.5), PptxInches(1.4)).table
            table.cell(0, 0).text = "Metric"
            table.cell(0, 1).text = "Value"
            table.cell(1, 0).text = "Coverage"
            table.cell(1, 1).text = "95"
        if with_image:
            slide.shapes.add_picture(str(self.image_file("pptx-embedded.png")), PptxInches(5.4), PptxInches(0.8), width=PptxInches(1.2))
        presentation.save(path)
        return path

    def xlsx_file(self, name="sample.xlsx"):
        path = Path(self.temp_media_dir) / name
        workbook = openpyxl.Workbook()
        ws = workbook.active
        ws.title = "Metrics"
        ws.append(["Player", "Runs"])
        ws.append(["Virat", 120])
        ws.append(["Rohit", 98])
        workbook.save(path)
        return path

    def csv_file(self, name="sample.csv"):
        path = Path(self.temp_media_dir) / name
        path.write_text("Name,Score\nVirat,120\nRohit,98\n", encoding="utf-8")
        return path

    def txt_file(self, name="sample.txt", content="Line 1\nLine 2\nLine 3\n"):
        path = Path(self.temp_media_dir) / name
        path.write_text(content, encoding="utf-8")
        return path

    def md_file(self, name="sample.md"):
        path = Path(self.temp_media_dir) / name
        path.write_text("# Overview\nThis is a markdown section.\n## Details\nSecond section.\n", encoding="utf-8")
        return path

    def svg_file(self, name="sample.svg"):
        path = Path(self.temp_media_dir) / name
        path.write_text(
            '<svg xmlns="http://www.w3.org/2000/svg" width="100" height="50"><text x="10" y="25">TechnoChat</text></svg>',
            encoding="utf-8",
        )
        return path


class AuthenticationTests(BaseTechnoChatTestCase):
    def test_tc_auth_01_valid_contributor_login_creates_session_and_redirects_home(self):
        user, _ = self.create_contributor()
        response = self.client.post(reverse("login"), {"email": user.email, "password": "StrongPass1!"})
        self.assertRedirects(response, reverse("home"))
        self.assertEqual(self.client.session["user_id"], user.id)

    def test_tc_auth_02_and_03_invalid_or_unknown_login_shows_error(self):
        user, _ = self.create_contributor()
        scenarios = [
            ("TC-AUTH-02", user.email, "WrongPass1!"),
            ("TC-AUTH-03", "missing@technostacks.com", "StrongPass1!"),
        ]
        for case_id, email, password in scenarios:
            with self.subTest(case_id=case_id):
                response = self.client.post(reverse("login"), {"email": email, "password": password})
                self.assertEqual(response.status_code, 200)
                self.assertContains(response, "Invalid email or password.")
                self.assertNotIn("user_id", self.client.session)

    def test_tc_auth_missing_credentials_show_validation_error(self):
        response = self.client.post(reverse("login"), {"email": "", "password": ""})
        self.assertEqual(response.status_code, 200)
        self.assertContains(response, "Please enter both email and password.")

    def test_tc_auth_04_first_login_redirects_to_profile(self):
        user = User.objects.create(email="new@technostacks.com", password="StrongPass1!", profile_completed=False)
        response = self.client.post(reverse("login"), {"email": user.email, "password": "StrongPass1!"})
        self.assertRedirects(response, reverse("profile"))
        self.assertTrue(UserProfile.objects.filter(user=user).exists())

    def test_tc_auth_05_profile_completion_saves_and_redirects_home(self):
        user, profile = self.create_contributor(profile_complete=False, username=None)
        self.set_contributor_session(user)
        response = self.client.post(
            reverse("profile"),
            {
                "first_name": "Raj",
                "surname": "Patel",
                "username": "raj.qa",
                "position_at_technostacks": "QA Engineer",
                "team": "QA",
            },
        )
        self.assertRedirects(response, reverse("home"))
        profile.refresh_from_db()
        user.refresh_from_db()
        self.assertTrue(profile.is_profile_complete)
        self.assertTrue(user.profile_completed)

    def test_tc_auth_06_and_07_session_persistence_and_logout_flush(self):
        user, _ = self.create_contributor()
        login_response = self.client.post(reverse("login"), {"email": user.email, "password": "StrongPass1!"})
        self.assertRedirects(login_response, reverse("home"))
        files_response = self.client.get(reverse("file_list"))
        self.assertEqual(files_response.status_code, 200)
        logout_response = self.client.get(reverse("logout"))
        self.assertRedirects(logout_response, reverse("login"))
        redirected = self.client.get(reverse("file_list"))
        self.assertRedirects(redirected, reverse("login"))

    def test_tc_auth_08_admin_login_with_valid_credentials_redirects_dashboard(self):
        admin, _ = self.create_admin_user()
        response = self.client.post(reverse("admin_login"), {"email": admin.email, "password": "StrongPass1!"})
        self.assertEqual(response.status_code, 302)
        self.assertEqual(response["Location"], reverse("admin_dashboard"))
        self.assertEqual(self.client.session["tc_admin_id"], admin.id)

    def test_tc_auth_09_admin_non_staff_login_is_blocked(self):
        admin, _ = self.create_admin_user(email="viewer@technostacks.com", is_staff=False)
        response = self.client.post(reverse("admin_login"), {"email": admin.email, "password": "StrongPass1!"})
        self.assertEqual(response.status_code, 200)
        self.assertContains(response, "You do not have admin access.")

    def test_tc_auth_10_admin_first_login_redirects_to_admin_profile(self):
        admin, profile = self.create_admin_user(email="newadmin@technostacks.com", profile_complete=False, username=None)
        response = self.client.post(reverse("admin_login"), {"email": admin.email, "password": "StrongPass1!"})
        self.assertRedirects(response, reverse("admin_profile"))
        profile.refresh_from_db()
        self.assertFalse(profile.is_profile_complete)

    def test_admin_logout_clears_admin_session_and_redirects_to_admin_login(self):
        admin, _ = self.create_admin_user()
        self.set_admin_session(admin)
        response = self.client.get(reverse("admin_logout"))
        self.assertRedirects(response, reverse("admin_login"))
        self.assertNotIn("tc_admin_id", self.client.session)


class FileUploadAndExtractionTests(BaseTechnoChatTestCase):
    def test_upload_requires_authenticated_contributor(self):
        upload = SimpleUploadedFile("report.pdf", b"dummy-content", content_type="application/pdf")
        response = self.client.post(reverse("upload_file"), {"file": upload})
        self.assertRedirects(response, reverse("login"))

    def test_tc_file_01_to_12_supported_upload_types_save_and_complete(self):
        user, _ = self.create_contributor()
        self.set_contributor_session(user)

        uploads = [
            ("TC-FILE-01", "report.pdf", "application/pdf", FileType.PDF),
            ("TC-FILE-02", "notes.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document", FileType.DOC),
            ("TC-FILE-03", "slides.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation", FileType.POWER),
            ("TC-FILE-04", "data.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", FileType.EXCEL),
            ("TC-FILE-05", "data.csv", "text/csv", FileType.CSV),
            ("TC-FILE-06", "notes.txt", "text/plain", FileType.TXT),
            ("TC-FILE-07", "readme.md", "text/markdown", FileType.MD),
            ("TC-FILE-08", "image.png", "image/png", FileType.IMAGE),
            ("TC-FILE-09", "image.jpg", "image/jpeg", FileType.IMAGE),
            ("TC-FILE-10", "image.jpeg", "image/jpeg", FileType.IMAGE),
            ("TC-FILE-11", "image.webp", "image/webp", FileType.IMAGE),
            ("TC-FILE-12", "image.svg", "image/svg+xml", FileType.IMAGE),
        ]

        def complete_embedding(file_object):
            file_object.embedding_status = FileProcessingStatus.COMPLETED
            file_object.save(update_fields=["embedding_status"])

        for case_id, name, content_type, expected_type in uploads:
            with self.subTest(case_id=case_id, name=name):
                File.objects.all().delete()
                upload = SimpleUploadedFile(name, b"dummy-content", content_type=content_type)
                with patch("backoffice_engine.views.embed_file_and_upsert", side_effect=complete_embedding) as embed_mock:
                    response = self.client.post(
                        reverse("upload_file"),
                        {"file": upload},
                        HTTP_X_REQUESTED_WITH="XMLHttpRequest",
                    )
                payload = response.json()
                self.assertTrue(payload["success"])
                saved = File.objects.get()
                self.assertEqual(saved.user, user)
                self.assertEqual(saved.original_filename, name)
                self.assertEqual(saved.file_type, expected_type)
                self.assertEqual(saved.embedding_status, FileProcessingStatus.COMPLETED)
                embed_mock.assert_called_once()

    def test_tc_file_13_to_15_upload_validation_errors(self):
        user, _ = self.create_contributor()
        self.set_contributor_session(user)
        huge_file = SimpleUploadedFile("large.pdf", b"x" * ((10 * 1024 * 1024) + 1), content_type="application/pdf")
        unsupported = SimpleUploadedFile("archive.zip", b"zip", content_type="application/zip")
        one = SimpleUploadedFile("one.txt", b"a", content_type="text/plain")
        two = SimpleUploadedFile("two.txt", b"b", content_type="text/plain")

        cases = [
            ("TC-FILE-13", {"file": huge_file}, "File too large."),
            ("TC-FILE-14", {"file": unsupported}, "File type not supported."),
            ("TC-FILE-15", {"file": [one, two]}, "Please upload one file at a time."),
        ]
        for case_id, payload, expected in cases:
            with self.subTest(case_id=case_id):
                response = self.client.post(reverse("upload_file"), payload, HTTP_X_REQUESTED_WITH="XMLHttpRequest")
                body = response.json()
                self.assertFalse(body["success"])
                self.assertIn(expected, body["error"])
                self.assertEqual(File.objects.count(), 0)

    @patch("backoffice_engine.document_reader._call_vlm_embedded", return_value="Embedded chart text")
    def test_tc_file_16_and_17_document_extractors_cover_images_tables_and_embedded_content(self, _mock_vlm):
        pdf_segments = extract_file_text(str(self.pdf_file(include_image=True)))
        self.assertIn("PDF Document Statistics", pdf_segments[0]["text"])
        self.assertTrue(any("Embedded chart text" in segment["text"] for segment in pdf_segments))

        docx_segments = extract_file_text(str(self.docx_file(with_table=True, with_image=True)))
        self.assertIn("Word Document Statistics", docx_segments[0]["text"])
        self.assertTrue(any("[Table:" in segment["text"] for segment in docx_segments))
        self.assertTrue(any("Embedded chart text" in segment["text"] for segment in docx_segments))

        pptx_segments = extract_file_text(str(self.pptx_file(with_table=True, with_image=True)))
        self.assertIn("PowerPoint Presentation Statistics", pptx_segments[0]["text"])
        self.assertTrue(any("[Table:" in segment["text"] for segment in pptx_segments))
        self.assertTrue(any("Embedded chart text" in segment["text"] for segment in pptx_segments))

        xlsx_segments = extract_file_text(str(self.xlsx_file()))
        self.assertIn("Excel File Statistics", xlsx_segments[0]["text"])
        self.assertTrue(any("[Sheet: Metrics]" in segment["text"] for segment in xlsx_segments))

        csv_segments = extract_file_text(str(self.csv_file()))
        self.assertIn("CSV File Statistical Summary", csv_segments[0]["text"])
        self.assertTrue(any("Row 2:" in segment["text"] for segment in csv_segments))

        txt_segments = extract_file_text(str(self.txt_file()))
        self.assertIn("Text File Statistics", txt_segments[0]["text"])
        self.assertTrue(any("Line 1" in segment["text"] for segment in txt_segments))

        md_segments = extract_file_text(str(self.md_file()))
        self.assertEqual(md_segments[0]["section_name"], "Document Statistics")
        self.assertTrue(any(segment["section_name"] == "Overview" for segment in md_segments))

        with patch("backoffice_engine.document_reader.VLMClient.describe_image_file", return_value="Standalone image description"):
            png_segments = extract_file_text(str(self.image_file("standalone.png")))
        self.assertIn("Image File Summary", png_segments[0]["text"])
        self.assertEqual(png_segments[1]["page_index"], 1)

        with patch("backoffice_engine.document_reader.VLMClient.describe_svg_file", return_value="Vector logo description"):
            svg_segments = extract_file_text(str(self.svg_file()))
        self.assertIn("Image File Summary", svg_segments[0]["text"])
        self.assertIn("Vector logo description", svg_segments[0]["text"])

    def test_embed_pipeline_builds_vectors_and_marks_completed(self):
        user, _ = self.create_contributor(email="pipeline@technostacks.com", username="pipeline.user")
        upload = SimpleUploadedFile("metrics.csv", b"Name,Score\nVirat,120\n", content_type="text/csv")
        file_obj = File.objects.create(user=user, file=upload, file_type=FileType.CSV, original_filename="metrics.csv")

        pinecone = MagicMock()
        pinecone.dense_text_embeddings.return_value = FakeDenseEmbeddings(2)
        pinecone.sparse_text_embeddings.return_value = FakeSparseEmbeddings(2)
        langchain = MagicMock()
        langchain.split_text.side_effect = lambda full_text: [full_text]

        with patch("backoffice_engine.ingestion_service.extract_file_text", return_value=[
            {"text": "CSV File Statistical Summary\nTotal rows: 2", "row_start": 1, "row_end": 2},
            {"text": "Row 2: Name: Virat | Score: 120", "row_start": 2, "row_end": 2},
        ]), patch("backoffice_engine.ingestion_service.PineconeClient", return_value=pinecone), patch(
            "backoffice_engine.ingestion_service.LangchainClient",
            return_value=langchain,
        ):
            embed_file_and_upsert(file_obj)

        file_obj.refresh_from_db()
        self.assertEqual(file_obj.embedding_status, FileProcessingStatus.COMPLETED)
        self.assertTrue(pinecone.upsert_file_data.called)
        vectors = pinecone.upsert_file_data.call_args.kwargs["vectors"]
        self.assertEqual(len(vectors), 2)
        self.assertEqual(vectors[0]["metadata"]["document_id"], file_obj.id)

    def test_tc_file_18_upload_empty_file_returns_no_text_error(self):
        user, _ = self.create_contributor(email="empty@technostacks.com", username="empty.user")
        self.set_contributor_session(user)
        upload = SimpleUploadedFile("empty.txt", b"", content_type="text/plain")
        response = self.client.post(
            reverse("upload_file"),
            {"file": upload},
            HTTP_X_REQUESTED_WITH="XMLHttpRequest",
        )
        payload = response.json()
        self.assertFalse(payload["success"])
        self.assertIn("The submitted file is empty.", payload["error"])
        self.assertEqual(File.objects.count(), 0)


class RagChatTests(BaseTechnoChatTestCase):
    def sample_chunks(self):
        return [
            {
                "file_id": 1,
                "file_name": "report.pdf",
                "file_type": "pdf",
                "normalized_file_type": "pdf",
                "page_index": 2,
                "score": 0.91,
                "chunk_index": 0,
                "text": "Virat Kohli scored 120 runs and Rohit Sharma scored 98 runs.",
            },
            {
                "file_id": 1,
                "file_name": "report.pdf",
                "file_type": "pdf",
                "normalized_file_type": "pdf",
                "page_index": 3,
                "score": 0.84,
                "chunk_index": 1,
                "text": "MS Dhoni has 3 centuries and 10 fifties in the file.",
            },
        ]

    def test_tc_rag_01_to_06_models_return_contextual_answers_with_sources(self):
        models = [
            "Gemini 2.5 Pro",
            "Gemini 2.5 Flash",
            "Gemini 2.5 Flash-Lite",
            "Llama 3.3 70B",
            "Llama 3.1 8B",
            "GPT OSS 120B",
        ]
        for model_name in models:
            with self.subTest(model=model_name):
                llm = FakeLLM(['{"answer": "Virat Kohli scored 120 runs.", "used_context_ids": [1]}'])
                with patch("backoffice_engine.chat_service.ChatPromptTemplate.from_messages", return_value=FakePromptPipe()), patch(
                    "backoffice_engine.chat_service._select_llm",
                    return_value=llm,
                ), patch(
                    "backoffice_engine.chat_service.retrieve_query_variations",
                    return_value=self.sample_chunks(),
                ), patch("backoffice_engine.chat_service.try_build_structured_answer", return_value=None):
                    result = chat_service.build_chat_prompt("Who scored 120 runs?", [1], [], model_name)
                self.assertIn("120 runs", result["answer"])
                self.assertEqual(result["sources"][0]["file_name"], "report.pdf")
                self.assertFalse(result["is_greeting"])

    def test_tc_rag_07_to_11_summary_list_exact_numeric_and_multi_question_behaviour(self):
        chunks = self.sample_chunks()

        summary_llm = FakeLLM(['{"answer": "Virat Kohli and MS Dhoni are covered in the document.", "used_context_ids": [1, 2]}'])
        with patch("backoffice_engine.chat_service.ChatPromptTemplate.from_messages", return_value=FakePromptPipe()), patch(
            "backoffice_engine.chat_service._select_llm",
            return_value=summary_llm,
        ), patch(
            "backoffice_engine.chat_service.retrieve_query_variations",
            return_value=chunks,
        ), patch("backoffice_engine.chat_service.try_build_structured_answer", return_value=None):
            summary = chat_service.build_chat_prompt("Give me summary of this file", [1], [], "Gemini 2.5 Pro")
        self.assertTrue(summary["is_summary"])
        self.assertIn("Virat Kohli", summary["answer"])

        list_llm = FakeLLM(['{"answer": "Virat Kohli\\nMS Dhoni", "used_context_ids": [1, 2]}'])
        with patch("backoffice_engine.chat_service.ChatPromptTemplate.from_messages", return_value=FakePromptPipe()), patch(
            "backoffice_engine.chat_service._select_llm",
            return_value=list_llm,
        ), patch(
            "backoffice_engine.chat_service.retrieve_query_variations",
            return_value=chunks,
        ), patch("backoffice_engine.chat_service.try_build_structured_answer", return_value=None):
            list_result = chat_service.build_chat_prompt("List all players", [1], [], "Gemini 2.5 Flash")
        self.assertIn("MS Dhoni", list_result["answer"])

        with patch("backoffice_engine.chat_service.retrieve_query_variations", return_value=chunks), patch(
            "backoffice_engine.chat_service._select_llm"
        ), patch("backoffice_engine.chat_service.try_build_structured_answer", return_value=None):
            exact = chat_service.build_chat_prompt("Give the word for word text", [1], [], "Llama 3.3 70B")
        self.assertEqual(exact["answer"], chunks[0]["text"])

        numeric_llm = FakeLLM(['{"answer": "Virat Kohli", "used_context_ids": [1]}'])
        with patch("backoffice_engine.chat_service.ChatPromptTemplate.from_messages", return_value=FakePromptPipe()), patch(
            "backoffice_engine.chat_service._select_llm",
            return_value=numeric_llm,
        ), patch(
            "backoffice_engine.chat_service.retrieve_query_variations",
            return_value=chunks,
        ), patch("backoffice_engine.chat_service.try_build_structured_answer", return_value=None):
            numeric = chat_service.build_chat_prompt("List players greater than 100 runs", [1], [], "GPT OSS 120B")
        self.assertEqual(numeric["answer"], "Virat Kohli")

        multi_llm = FakeLLM([
            '{"answer": "Virat Kohli scored 120 runs.", "used_context_ids": [1]}',
            '{"answer": "MS Dhoni has 3 centuries.", "used_context_ids": [2]}',
        ])
        with patch("backoffice_engine.chat_service.ChatPromptTemplate.from_messages", return_value=FakePromptPipe()), patch(
            "backoffice_engine.chat_service._select_llm",
            return_value=multi_llm,
        ), patch(
            "backoffice_engine.chat_service.retrieve_query_variations",
            return_value=chunks,
        ), patch("backoffice_engine.chat_service.try_build_structured_answer", return_value=None):
            multi = chat_service.build_chat_prompt("Who scored 120 runs? And how many centuries does MS Dhoni have?", [1], [], "Gemini 2.5 Flash")
        self.assertIn("Virat Kohli scored 120 runs.", multi["answer"])
        self.assertIn("MS Dhoni has 3 centuries.", multi["answer"])

    def test_tc_rag_12_and_13_structured_page_count_and_vague_follow_up_resolution(self):
        with patch("backoffice_engine.chat_service.try_build_structured_answer", return_value={
            "answer": "The PDF has 2 page(s).",
            "sources": [{"file_id": 1, "file_name": "report.pdf", "file_type": "pdf", "page_index": 1}],
        }):
            result = chat_service.build_chat_prompt("What is the page count?", [1], [], "Gemini 2.5 Pro")
        self.assertEqual(result["answer"], "The PDF has 2 page(s).")
        self.assertEqual(result["sources"][0]["page_index"], 1)

        llm = FakeLLM(['{"answer": "It refers to cricket statistics and player performance.", "used_context_ids": [1]}'])
        with patch("backoffice_engine.chat_service.ChatPromptTemplate.from_messages", return_value=FakePromptPipe()), patch(
            "backoffice_engine.chat_service._select_llm",
            return_value=llm,
        ), patch(
            "backoffice_engine.chat_service.retrieve_query_variations",
            return_value=self.sample_chunks(),
        ), patch("backoffice_engine.chat_service.try_build_structured_answer", return_value=None):
            follow_up = chat_service.build_chat_prompt(
                "What about it?",
                [1],
                [SimpleNamespace(question="Tell me about cricket statistics", answer="...")],
                "Gemini 2.5 Pro",
                conversation_state={"active_topic": "cricket statistics", "active_entities": ["MS Dhoni"]},
            )
        self.assertIn("cricket statistics", follow_up["resolved_query"])

    def test_tc_rag_14_returns_empty_response_when_no_context_found(self):
        with patch("backoffice_engine.chat_service.retrieve_query_variations", return_value=[]), patch(
            "backoffice_engine.chat_service.try_build_structured_answer",
            return_value=None,
        ), patch("backoffice_engine.chat_service._select_llm"):
            result = chat_service.build_chat_prompt("Question with no match", [1], [], "Gemini 2.5 Pro")
        self.assertEqual(result["answer"], chat_service.EMPTY_RAG_RESPONSE)
        self.assertEqual(result["sources"], [])

    def test_tc_rag_abusive_non_question_returns_respectful_response(self):
        result = chat_service.build_chat_prompt("idiot", [1], [], "Gemini 2.5 Pro")
        self.assertEqual(result["answer"], chat_service.RESPECTFUL_RESPONSE)
        self.assertEqual(result["sources"], [])

    def test_tc_rag_15_to_17_source_viewer_endpoints_handle_pdf_pptx_and_text(self):
        user, _ = self.create_contributor(email="viewer@technostacks.com", username="viewer.user")
        self.set_contributor_session(user)

        pdf = File.objects.create(
            user=user,
            file=SimpleUploadedFile("report.pdf", self.pdf_file().read_bytes(), content_type="application/pdf"),
            file_type=FileType.PDF,
            original_filename="report.pdf",
            embedding_status=FileProcessingStatus.COMPLETED,
        )
        pptx = File.objects.create(
            user=user,
            file=SimpleUploadedFile("slides.pptx", self.pptx_file().read_bytes(), content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"),
            file_type=FileType.POWER,
            original_filename="slides.pptx",
            embedding_status=FileProcessingStatus.COMPLETED,
        )
        txt = File.objects.create(
            user=user,
            file=SimpleUploadedFile("notes.txt", b"Line A\nLine B\nLine C\n", content_type="text/plain"),
            file_type=FileType.TXT,
            original_filename="notes.txt",
            embedding_status=FileProcessingStatus.COMPLETED,
        )

        pdf_response = self.client.get(reverse("page_render"), {"file_id": pdf.id, "file_type": "pdf", "page_index": 1})
        self.assertTrue(pdf_response.json()["success"])
        self.assertEqual(pdf_response.json()["source_type"], "page")

        pptx_response = self.client.get(reverse("page_render"), {"file_id": pptx.id, "file_type": "pptx", "slide_index": 1})
        self.assertTrue(pptx_response.json()["success"])
        self.assertEqual(pptx_response.json()["source_type"], "page")

        txt_response = self.client.get(reverse("page_render"), {"file_id": txt.id, "file_type": "txt", "line_start": 1, "line_end": 2})
        self.assertTrue(txt_response.json()["success"])
        self.assertEqual(txt_response.json()["source_type"], "text")
        self.assertIn("Line A", txt_response.json()["content_text"])


class AIAssistantTests(BaseTechnoChatTestCase):
    def test_tc_ai_01_02_06_07_general_and_coding_questions_work_across_models(self):
        models = ["Gemini 2.5 Pro", "Llama 3.3 70B"]
        for model in models:
            with self.subTest(model=model):
                llm = FakeLLM([
                    '{"answer": "Artificial intelligence is the simulation of human intelligence by machines."}',
                    '{"answer": "```python\\nprint(\\"hello\\")\\n```"}',
                ])
                with patch("backoffice_engine.ai_assistant_service.ChatPromptTemplate.from_messages", return_value=FakePromptPipe()), patch(
                    "backoffice_engine.ai_assistant_service._select_llm",
                    return_value=llm,
                ):
                    general = ai_assistant_service.build_ai_assistant_prompt("What is AI?", [], model)
                    coding = ai_assistant_service.build_ai_assistant_prompt("Write Python code to print hello", [], model)
                self.assertEqual(general["sources"], [])
                self.assertIn("Artificial intelligence", general["answer"])
                self.assertIn("python", coding["answer"].lower())

    def test_tc_ai_03_and_05_greeting_and_abuse_are_handled_without_model_call(self):
        greeting = ai_assistant_service.build_ai_assistant_prompt("Hi", [], "Gemini 2.5 Pro")
        self.assertTrue(greeting["is_greeting"])
        self.assertEqual(greeting["answer"], "Hi! How can I help you today?")

        abusive = ai_assistant_service.build_ai_assistant_prompt("idiot", [], "Gemini 2.5 Pro")
        self.assertEqual(abusive["answer"], ai_assistant_service.RESPECTFUL_RESPONSE)

    def test_tc_ai_04_follow_up_uses_history_and_tc_ai_08_handles_multi_part_questions(self):
        llm = FakeLLM([
            '{"answer": "It continues the previous discussion."}',
            '{"answer": "Part one answer."}',
            '{"answer": "Part two answer."}',
        ])
        with patch("backoffice_engine.ai_assistant_service.ChatPromptTemplate.from_messages", return_value=FakePromptPipe()), patch(
            "backoffice_engine.ai_assistant_service._select_llm",
            return_value=llm,
        ):
            follow_up = ai_assistant_service.build_ai_assistant_prompt(
                "What about it?",
                [SimpleNamespace(question="Tell me about AI agents", answer="AI agents can plan.")],
                "Gemini 2.5 Flash",
                conversation_state={"active_topic": "AI agents", "active_entities": ["AI agents"]},
            )
            multi = ai_assistant_service.build_ai_assistant_prompt(
                "What is Python? And what is Django?",
                [],
                "Gemini 2.5 Flash",
            )
        self.assertIn("artificial intelligence agents", follow_up["resolved_query"].lower())
        self.assertIn("Part one answer.", multi["answer"])
        self.assertIn("Part two answer.", multi["answer"])
        self.assertTrue(llm.calls[0]["chat_history"])


class WebSearchTests(BaseTechnoChatTestCase):
    def test_tc_web_01_to_03_news_and_live_data_queries_return_sources(self):
        llm = FakeLLM([
            '{"answer": "Current news summary.", "used_result_ids": [1, 2]}',
            '{"answer": "Weather is 29C.", "used_result_ids": [1]}',
        ])
        results = [
            {"title": "News One", "link": "https://news.example/1", "snippet": "Snippet 1"},
            {"title": "News Two", "link": "https://news.example/2", "snippet": "Snippet 2"},
        ]
        with patch("backoffice_engine.web_search_service._select_llm", return_value=llm), patch(
            "backoffice_engine.web_search_service.SerperClient.search",
            return_value=results,
        ):
            news = web_search_service.build_web_search_prompt("Latest AI news", "Gemini 2.5 Pro", [])
            weather = web_search_service.build_web_search_prompt("Current weather in Ahmedabad", "Gemini 2.5 Pro", [])
        self.assertEqual(len(news["sources"]), 2)
        self.assertEqual(news["sources"][0]["link"], "https://news.example/1")
        self.assertIn("Weather is 29C.", weather["answer"])

    def test_tc_web_04_and_05_greeting_and_context_resolution(self):
        greeting = web_search_service.build_web_search_prompt("Hi", "Gemini 2.5 Pro", [])
        self.assertTrue(greeting["is_greeting"])
        self.assertEqual(greeting["sources"], [])

        llm = FakeLLM(['{"answer": "Topic follow-up answer.", "used_result_ids": [1]}'])
        with patch("backoffice_engine.web_search_service._select_llm", return_value=llm), patch(
            "backoffice_engine.web_search_service.SerperClient.search",
            return_value=[{"title": "Topic", "link": "https://topic.example", "snippet": "Topic snippet"}],
        ) as search_mock:
            web_search_service.build_web_search_prompt(
                "What about it?",
                "Gemini 2.5 Pro",
                [SimpleNamespace(question="Tell me about Indian stock market", answer="...")],
                conversation_state={"active_topic": "Indian stock market", "active_entities": ["NSE"]},
            )
        self.assertIn("Indian stock market", search_mock.call_args.args[0])

    def test_tc_web_06_no_results_raises_graceful_error(self):
        llm = FakeLLM(['{"answer": "unused", "used_result_ids": []}'])
        with patch("backoffice_engine.web_search_service._select_llm", return_value=llm), patch(
            "backoffice_engine.web_search_service.SerperClient.search",
            return_value=[],
        ):
            with self.assertRaises(WebSearchError):
                web_search_service.build_web_search_prompt("No results query", "Gemini 2.5 Pro", [])

    def test_tc_web_abusive_non_question_returns_respectful_response(self):
        result = web_search_service.build_web_search_prompt("idiot", "Gemini 2.5 Pro", [])
        self.assertEqual(result["answer"], web_search_service.RESPECTFUL_RESPONSE)
        self.assertEqual(result["sources"], [])


class ImageGenerationTests(BaseTechnoChatTestCase):
    def test_tc_img_01_and_02_text_to_image_and_image_edit(self):
        request = MagicMock()
        request.build_absolute_uri.side_effect = lambda url: f"http://testserver{url}"
        upload = SimpleUploadedFile("photo.png", b"img", content_type="image/png")

        client = MagicMock()
        client.api_key = "key"
        client.text_model = "gpt-image"
        client.edit_model = "gpt-image-edit"
        client.text_to_image.return_value = ["https://img.example/generated.png"]
        client.image_to_image.return_value = ["https://img.example/edited.png"]

        with patch("backoffice_engine.image_generation_service.KieImageClient", return_value=client):
            generated = image_generation_service.build_image_generation_prompt("Create a mountain scene", request)
        self.assertEqual(generated["image_urls"][0], "https://img.example/generated.png")
        self.assertEqual(generated["sources"][0]["kind"], "generated_image")

        with patch("backoffice_engine.image_generation_service.KieImageClient", return_value=client), patch(
            "backoffice_engine.image_generation_service.save_uploaded_image",
            return_value="/media/chat_inputs/photo.png",
        ), patch(
            "backoffice_engine.image_generation_service.uploaded_image_to_data_uri",
            return_value="data:image/png;base64,abc",
        ):
            edited = image_generation_service.build_image_generation_prompt("Add warm lighting", request, upload)
        self.assertEqual(edited["answer"], "Here is the edited image.")
        self.assertEqual(edited["sources"][0]["kind"], "uploaded_image")

    def test_tc_img_05_and_06_chat_send_view_rejects_wrong_mode_and_missing_prompt(self):
        user, _ = self.create_contributor(email="img@technostacks.com", username="img.user")
        session = ChatSession.objects.create(user=user, title="Image Session", session_type=SessionType.GENERAL_CHAT)
        self.set_contributor_session(user)
        image = SimpleUploadedFile("photo.png", b"img", content_type="image/png")

        wrong_mode = self.client.post(
            reverse("chat_send", args=[session.id]),
            {"query": "hello", "model_name": "Gemini 2.5 Pro", "chat_mode": "ai_assistant", "image": image},
        )
        self.assertEqual(wrong_mode.status_code, 400)
        self.assertIn("Select Create Image mode to upload an image.", wrong_mode.json()["error"])

        image = SimpleUploadedFile("photo.png", b"img", content_type="image/png")
        missing_prompt = self.client.post(
            reverse("chat_send", args=[session.id]),
            {"query": "", "model_name": "Gemini 2.5 Pro", "chat_mode": "image_generation", "image": image},
        )
        self.assertEqual(missing_prompt.status_code, 400)
        self.assertIn("Please enter a message before sending.", missing_prompt.json()["error"])

    def test_tc_img_08_timeout_is_converted_to_user_friendly_error(self):
        request = MagicMock()
        client = MagicMock()
        client.api_key = "key"
        client.text_model = "gpt-image"
        client.edit_model = "gpt-image-edit"
        client.text_to_image.side_effect = TimeoutError("slow")
        with patch("backoffice_engine.image_generation_service.KieImageClient", return_value=client):
            with self.assertRaises(ChatResponseError) as exc:
                image_generation_service.build_image_generation_prompt("Create art", request)
        self.assertIn("Image generation is taking longer", exc.exception.user_message)

    def test_tc_img_unconfigured_client_returns_configuration_error(self):
        request = MagicMock()
        client = MagicMock()
        client.api_key = ""
        client.text_model = ""
        client.edit_model = ""
        with patch("backoffice_engine.image_generation_service.KieImageClient", return_value=client):
            with self.assertRaises(ChatResponseError) as exc:
                image_generation_service.build_image_generation_prompt("Create art", request)
        self.assertIn("not configured properly", exc.exception.user_message)


class SessionAndChatManagementTests(BaseTechnoChatTestCase):
    def test_tc_sess_01_and_02_create_file_and_general_sessions(self):
        user, _ = self.create_contributor(email="session@technostacks.com", username="session.user")
        self.set_contributor_session(user)
        file_obj = File.objects.create(
            user=user,
            file=SimpleUploadedFile("report.pdf", b"pdf", content_type="application/pdf"),
            file_type=FileType.PDF,
            original_filename="report.pdf",
            embedding_status=FileProcessingStatus.COMPLETED,
        )

        file_response = self.client.post(
            reverse("create_session"),
            {"title": "Project Chat", "session_type": "chat_with_file", "file_ids": [file_obj.id]},
        )
        created = ChatSession.objects.get(title="Project Chat")
        self.assertRedirects(file_response, reverse("chat", args=[created.id]))
        self.assertEqual(list(created.files.values_list("id", flat=True)), [file_obj.id])

        general_response = self.client.post(
            reverse("create_session"),
            {"title": "General Talk", "session_type": "general_chat"},
        )
        general = ChatSession.objects.get(title="General Talk")
        self.assertRedirects(general_response, reverse("chat", args=[general.id]))
        self.assertEqual(general.files.count(), 0)

    def test_tc_sess_03_to_05_session_sidebar_message_count_and_resume(self):
        user, _ = self.create_contributor(email="resume@technostacks.com", username="resume.user")
        self.set_contributor_session(user)
        session = ChatSession.objects.create(user=user, title="Project Chat", session_type=SessionType.GENERAL_CHAT)
        ChatMessage.objects.create(session=session, question="Q1", answer="A1", model_used="Gemini 2.5 Pro", chat_mode="web_search")
        response = self.client.get(reverse("chat", args=[session.id]))
        self.assertContains(response, "Project Chat")
        self.assertContains(response, 'id="count-')
        self.assertContains(response, "Q1")

        list_response = self.client.get(reverse("chat_list"))
        self.assertRedirects(list_response, reverse("chat", args=[session.id]))

    def test_tc_sess_06_to_08_chat_send_updates_mode_last_active_and_persists_message(self):
        user, _ = self.create_contributor(email="message@technostacks.com", username="message.user")
        self.set_contributor_session(user)
        session = ChatSession.objects.create(user=user, title="General Talk", session_type=SessionType.GENERAL_CHAT)

        with patch("backoffice_engine.views.build_ai_assistant_prompt", return_value={
            "answer": "AI answer",
            "sources": [],
            "is_greeting": False,
            "is_summary": False,
            "chat_mode": "ai_assistant",
            "resolved_query": "What is AI?",
        }):
            response = self.client.post(
                reverse("chat_send", args=[session.id]),
                data=json.dumps({"query": "What is AI?", "model_name": "Gemini 2.5 Pro", "chat_mode": "ai_assistant"}),
                content_type="application/json",
            )

        payload = response.json()
        self.assertTrue(payload["success"])
        self.assertEqual(payload["message"]["message_count"], 1)
        self.assertEqual(payload["message"]["chat_mode"], "ai_assistant")
        self.assertEqual(ChatMessage.objects.filter(session=session).count(), 1)
        saved = ChatMessage.objects.get(session=session)
        self.assertEqual(saved.sources, [])
        self.assertEqual(get_last_active_chat_session_id(self.client), session.id)

    def test_chat_send_requires_post_method(self):
        user, _ = self.create_contributor(email="method@technostacks.com", username="method.user")
        self.set_contributor_session(user)
        session = ChatSession.objects.create(user=user, title="General Talk", session_type=SessionType.GENERAL_CHAT)
        response = self.client.get(reverse("chat_send", args=[session.id]))
        self.assertEqual(response.status_code, 405)
        self.assertEqual(response.json()["error"], "POST required")

    def test_page_render_requires_auth_and_file_id(self):
        unauthenticated = self.client.get(reverse("page_render"))
        self.assertEqual(unauthenticated.status_code, 401)
        self.assertEqual(unauthenticated.json()["error"], "Not authenticated")

        user, _ = self.create_contributor(email="preview@technostacks.com", username="preview.user")
        self.set_contributor_session(user)
        missing_id = self.client.get(reverse("page_render"))
        self.assertEqual(missing_id.status_code, 400)
        self.assertEqual(missing_id.json()["error"], "file_id required")


class AdminPortalTests(BaseTechnoChatTestCase):
    def test_tc_adm_01_02_and_09_dashboard_and_lists_render(self):
        admin, _ = self.create_admin_user()
        other_admin, _ = self.create_admin_user(email="second@technostacks.com", username="second.admin")
        contributor, profile = self.create_contributor(email="contrib@technostacks.com", username="contrib.user")
        File.objects.create(
            user=contributor,
            file=SimpleUploadedFile("report.pdf", b"pdf", content_type="application/pdf"),
            file_type=FileType.PDF,
            original_filename="report.pdf",
            embedding_status=FileProcessingStatus.COMPLETED,
        )
        self.set_admin_session(admin)

        captured = {}
        with patch("backoffice_engine.admin_views.render", side_effect=fake_render_factory(captured)):
            admins_response = self.client.get(reverse("admin_dashboard"), {"section": "admins"})
        self.assertEqual(admins_response.status_code, 200)
        admin_emails = list(captured["context"]["records"].values_list("email", flat=True))
        self.assertIn(admin.email, admin_emails)
        self.assertIn(other_admin.email, admin_emails)

        captured = {}
        with patch("backoffice_engine.admin_views.render", side_effect=fake_render_factory(captured)):
            profiles_response = self.client.get(reverse("admin_dashboard"), {"section": "profiles"})
        self.assertEqual(profiles_response.status_code, 200)
        profile_records = list(captured["context"]["records"])
        self.assertTrue(any(item.pk == profile.pk and item.team == profile.team for item in profile_records))

    def test_tc_adm_03_edit_admin_email_logs_change(self):
        admin, _ = self.create_admin_user()
        target, _ = self.create_admin_user(email="editable@technostacks.com", username="editable.admin")
        self.set_admin_session(admin)
        response = self.client.post(
            reverse("admin_dashboard") + f"?section=admins&edit_id={target.id}",
            {"action": "edit_save", "email": "updated@technostacks.com", "profile_completed": "on"},
        )
        self.assertEqual(response.status_code, 302)
        target.refresh_from_db()
        self.assertEqual(target.email, "updated@technostacks.com")
        self.assertTrue(LogEntry.objects.filter(object_id=str(target.pk), action_flag=CHANGE).exists())

    def test_tc_adm_04_delete_single_admin_cannot_delete_self_but_can_delete_other(self):
        admin, _ = self.create_admin_user()
        other, _ = self.create_admin_user(email="remove@technostacks.com", username="remove.admin")
        self.set_admin_session(admin)

        self.client.post(
            reverse("admin_dashboard") + f"?section=admins&edit_id={admin.id}",
            {"action": "delete_single"},
        )
        self.assertTrue(AdminUser.objects.filter(pk=admin.pk).exists())

        self.client.post(
            reverse("admin_dashboard") + f"?section=admins&edit_id={other.id}",
            {"action": "delete_single"},
        )
        self.assertFalse(AdminUser.objects.filter(pk=other.pk).exists())

    def test_tc_adm_05_bulk_delete_contributors(self):
        admin, _ = self.create_admin_user()
        user_one, _ = self.create_contributor(email="bulk1@technostacks.com", username="bulk.one")
        user_two, _ = self.create_contributor(email="bulk2@technostacks.com", username="bulk.two")
        self.set_admin_session(admin)
        response = self.client.post(
            reverse("admin_dashboard") + "?section=contributors",
            {"action": "delete_bulk", "record_ids": [user_one.id, user_two.id]},
        )
        self.assertEqual(response.status_code, 302)
        self.assertFalse(User.objects.filter(pk__in=[user_one.id, user_two.id]).exists())

    def test_tc_adm_06_07_08_and_16_new_contributor_creation_validation_and_popup(self):
        admin, _ = self.create_admin_user()
        self.set_admin_session(admin)

        invalid_domain = self.client.post(
            reverse("admin_new_contributor"),
            {"email": "person@example.com", "password": "StrongPass1!", "position": "QA", "team": "QA"},
        )
        self.assertEqual(invalid_domain.status_code, 302)
        self.assertEqual(invalid_domain["Location"], reverse("admin_dashboard") + "?section=contributors")
        invalid_messages = [message.message for message in get_messages(invalid_domain.wsgi_request)]
        self.assertTrue(any("Only @technostacks.com email addresses are allowed." in message for message in invalid_messages))

        weak_password = self.client.post(
            reverse("admin_new_contributor"),
            {"email": "weak@technostacks.com", "password": "weak", "position": "QA", "team": "QA"},
        )
        self.assertEqual(weak_password.status_code, 302)
        self.assertEqual(weak_password["Location"], reverse("admin_dashboard") + "?section=contributors")
        weak_messages = [message.message for message in get_messages(weak_password.wsgi_request)]
        self.assertTrue(any("Password must be at least 8 characters." in message for message in weak_messages))

        success = self.client.post(
            reverse("admin_new_contributor"),
            {"email": "newcontrib@technostacks.com", "password": "StrongPass1!", "position": "QA", "team": "QA"},
        )
        self.assertEqual(success.status_code, 302)
        self.assertEqual(success["Location"], reverse("admin_dashboard") + "?section=contributors&contributor_added=1")
        self.assertTrue(User.objects.filter(email="newcontrib@technostacks.com").exists())
        self.assertTrue(UserProfile.objects.filter(user__email="newcontrib@technostacks.com").exists())

        captured = {}
        with patch("backoffice_engine.admin_views.render", side_effect=fake_render_factory(captured)):
            popup = self.client.get(reverse("admin_dashboard"), {"section": "contributors", "contributor_added": "1"})
        self.assertEqual(popup.status_code, 200)
        self.assertTrue(captured["context"]["show_contributor_success_popup"])

    def test_tc_adm_10_11_12_13_14_15_history_file_session_and_message_sections(self):
        admin, _ = self.create_admin_user()
        contributor, profile = self.create_contributor(email="full@technostacks.com", username="full.user", team="QA")
        file_obj = File.objects.create(
            user=contributor,
            file=SimpleUploadedFile("report.pdf", b"pdf", content_type="application/pdf"),
            file_type=FileType.PDF,
            original_filename="report.pdf",
            embedding_status=FileProcessingStatus.PENDING,
        )
        session = ChatSession.objects.create(user=contributor, title="Project Chat", session_type=SessionType.CHAT_WITH_FILE)
        session.files.set([file_obj])
        ChatMessage.objects.create(
            session=session,
            question="What is AI?",
            answer="AI is...",
            model_used="Gemini 2.5 Pro",
            chat_mode="rag",
            sources=[{"file_name": "report.pdf", "page_index": 2}],
        )
        self.set_admin_session(admin)

        edit_profile = self.client.post(
            reverse("admin_dashboard") + f"?section=profiles&edit_id={profile.id}",
            {
                "action": "edit_save",
                "first_name": "Full",
                "surname": "User",
                "username": "full.user",
                "position_at_technostacks": "QA Lead",
                "team": "Core",
            },
        )
        self.assertEqual(edit_profile.status_code, 302)
        profile.refresh_from_db()
        self.assertEqual(profile.team, "Core")

        file_edit = self.client.post(
            reverse("admin_dashboard") + f"?section=files&edit_id={file_obj.id}",
            {"action": "edit_save", "original_filename": "report-final.pdf", "embedding_status": "completed"},
        )
        self.assertEqual(file_edit.status_code, 302)
        file_obj.refresh_from_db()
        self.assertEqual(file_obj.embedding_status, FileProcessingStatus.COMPLETED)

        captured = {}
        with patch("backoffice_engine.admin_views.render", side_effect=fake_render_factory(captured)):
            sessions_response = self.client.get(reverse("admin_dashboard"), {"section": "sessions"})
        self.assertEqual(sessions_response.status_code, 200)
        session_records = list(captured["context"]["records"])
        self.assertTrue(any(item.pk == session.pk and item.title == "Project Chat" for item in session_records))

        captured = {}
        with patch("backoffice_engine.admin_views.render", side_effect=fake_render_factory(captured)):
            messages_response = self.client.get(reverse("admin_dashboard"), {"section": "messages"})
        self.assertEqual(messages_response.status_code, 200)
        message_records = list(captured["context"]["records"])
        self.assertTrue(any(item.question == "What is AI?" and item.chat_mode == "rag" for item in message_records))

        captured = {}
        with patch(
            "backoffice_engine.admin_views._history_entries_for_object",
            return_value=[
                {
                    "label": "Updated",
                    "tone": "info",
                    "time": profile.updated_at,
                    "user": admin.email,
                    "message": "Profile updated.",
                }
            ],
        ), patch("backoffice_engine.admin_views.render", side_effect=fake_render_factory(captured)):
            history_response = self.client.get(reverse("admin_dashboard"), {"section": "profiles", "edit_id": profile.id, "history": "1"})
        self.assertEqual(history_response.status_code, 200)
        self.assertTrue(captured["context"]["show_history"])
        self.assertEqual(captured["context"]["history_button_label"], "Hide History")
        self.assertGreaterEqual(len(captured["context"]["history_entries"]), 1)


class UINavigationTests(BaseTechnoChatTestCase):
    def test_tc_ui_01_to_15_templates_route_and_render_expected_controls(self):
        user, _ = self.create_contributor(email="ui@technostacks.com", username="ui.user")
        self.set_contributor_session(user)
        file_obj = File.objects.create(
            user=user,
            file=SimpleUploadedFile("report.pdf", b"pdf", content_type="application/pdf"),
            file_type=FileType.PDF,
            original_filename="report.pdf",
            embedding_status=FileProcessingStatus.COMPLETED,
        )
        session = ChatSession.objects.create(user=user, title="Project Chat", session_type=SessionType.CHAT_WITH_FILE)
        session.files.set([file_obj])
        ChatMessage.objects.create(
            session=session,
            question="What is AI?",
            answer="AI is...",
            model_used="Gemini 2.5 Pro",
            chat_mode="rag",
            sources=[{"file_name": "report.pdf", "file_id": file_obj.id, "file_type": "pdf", "page_index": 2}],
        )

        routes = [
            reverse("home"),
            reverse("file_list"),
            reverse("chat", args=[session.id]),
            reverse("about_us"),
            reverse("create_session"),
        ]
        for route in routes:
            with self.subTest(route=route):
                response = self.client.get(route)
                self.assertEqual(response.status_code, 200)

        chat_response = self.client.get(reverse("chat", args=[session.id]))
        self.assertContains(chat_response, "source-viewer-modal")
        self.assertContains(chat_response, "typing-wrapper")

        about_response = self.client.get(reverse("about_us"))
        self.assertContains(about_response, "Technology Stack")
        self.assertContains(about_response, "Django")

    def test_tc_ui_16_to_18_static_js_contains_expected_interaction_logic(self):
        base_js = (Path(settings.BASE_DIR) / "backoffice_engine" / "static" / "base.js").read_text(encoding="utf-8")
        upload_js = (Path(settings.BASE_DIR) / "backoffice_engine" / "static" / "upload.js").read_text(encoding="utf-8")
        chat_js = (Path(settings.BASE_DIR) / "backoffice_engine" / "static" / "chat.js").read_text(encoding="utf-8")
        create_session_js = (Path(settings.BASE_DIR) / "backoffice_engine" / "static" / "create_session.js").read_text(encoding="utf-8")
        admin_dash_js = (Path(settings.BASE_DIR) / "backoffice_engine" / "static" / "admin" / "admin_dashboard.js").read_text(encoding="utf-8")
        admin_dash_html = (Path(settings.BASE_DIR) / "backoffice_engine" / "templates" / "admin" / "admin_dashboard.html").read_text(encoding="utf-8")

        self.assertIn("profileDropdown.classList.toggle('show')", base_js)
        self.assertIn("localStorage.setItem('tc_theme', theme)", base_js)
        self.assertIn("zone.classList.toggle('active')", upload_js)
        self.assertIn("btnSend.disabled = chatInput.value.trim() === '' && !pendingImageFile;", chat_js)
        self.assertIn("Math.min(this.scrollHeight, 160)", chat_js)
        self.assertIn("if (event.key === 'Enter' && !event.shiftKey)", chat_js)
        self.assertIn("typingWrapper.style.display = 'flex';", chat_js)
        self.assertIn("modal.style.display = 'flex';", chat_js)
        self.assertIn("btn-select-all", create_session_js)
        self.assertIn("btn-deselect-all", create_session_js)
        self.assertIn("currentChatMode = btn.getAttribute('data-mode')", chat_js)
        self.assertIn("setupToggle('e-pw-toggle-contrib'", admin_dash_js)
        self.assertIn("m-r-len", admin_dash_html)
        self.assertIn("History", admin_dash_html)


class AdminModelAndHelperSmokeTests(BaseTechnoChatTestCase):
    def test_create_admin_helper_initializes_profile(self):
        admin = create_admin("freshadmin@technostacks.com", "StrongPass1!")
        self.assertTrue(AdminProfile.objects.filter(admin=admin).exists())

    def test_chat_list_uses_last_active_session_when_present(self):
        user, _ = self.create_contributor(email="last@technostacks.com", username="last.user")
        self.set_contributor_session(user)
        first = ChatSession.objects.create(user=user, title="First", session_type=SessionType.GENERAL_CHAT)
        second = ChatSession.objects.create(user=user, title="Second", session_type=SessionType.GENERAL_CHAT)
        self.client.get(reverse("chat", args=[second.id]))
        response = self.client.get(reverse("chat_list"))
        self.assertRedirects(response, reverse("chat", args=[second.id]))
        self.assertNotEqual(first.id, get_last_active_chat_session_id(self.client))

    def test_chat_list_without_sessions_redirects_to_create_session(self):
        user, _ = self.create_contributor(email="emptychat@technostacks.com", username="emptychat.user")
        self.set_contributor_session(user)
        response = self.client.get(reverse("chat_list"))
        self.assertRedirects(response, reverse("create_session"))


class StaticRegressionSmokeTests(SimpleTestCase):
    def test_templates_keep_password_eye_and_history_controls_present(self):
        base_dir = Path(settings.BASE_DIR) / "backoffice_engine"
        admin_dashboard = (base_dir / "templates" / "admin" / "admin_dashboard.html").read_text(encoding="utf-8")
        create_session = (base_dir / "templates" / "create_session.html").read_text(encoding="utf-8")
        chat_template = (base_dir / "templates" / "chat.html").read_text(encoding="utf-8")
        self.assertIn("e-pw-toggle-contrib", admin_dashboard)
        self.assertIn("m-pw-t1", admin_dashboard)
        self.assertIn("history-link", admin_dashboard)
        self.assertIn("btn-select-all", create_session)
        self.assertIn("source-viewer-modal", chat_template)
